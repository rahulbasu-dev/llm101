"""Generate a PPTX from the 12 Build Steps.

Run: python gen_pptx.py
Output: LLM101_Build_Steps.pptx
"""

import os
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

import build_viz

# ── Render diagrams ──
img_dir = tempfile.mkdtemp(prefix="llm101_pptx_")
diagram_paths = build_viz.render_all(img_dir)

IMG_MAP = {
    "config": diagram_paths.get("config"),
    "sliding": diagram_paths.get("sliding"),
    "block": diagram_paths.get("block"),
    "training": diagram_paths.get("training"),
    "kv_cache": diagram_paths.get("kv_cache"),
    "ui_layout": diagram_paths.get("ui_layout"),
    "test_matrix": diagram_paths.get("test_matrix"),
}

# ── Step data: (title, body, image_key_or_None) ──
STEPS = [
    ("1 \u00b7 Pin the hyperparameters \u2014 config.py",
     "Start here, not at model.py. A single dataclass owning every tunable "
     "keeps the rest of the code uncluttered and reproducibility mechanical.\n\n"
     "Non-obvious choice: vocab_size starts at 0. Every script that builds "
     "the model must set config.vocab_size = tokenizer.vocab_size before "
     "instantiation. This avoids hardcoding and keeps the tokenizer and model "
     "coupled through a single variable.\n\n"
     "Derived values live in @property, with assertions that catch "
     "misconfiguration at access time rather than during a forward pass.",
     "config"),

    ("2 \u00b7 Byte-level BPE tokenizer \u2014 tokenizer.py",
     "BPE from scratch in ~200 lines. Byte-level means the base vocab is "
     "always 256 (every possible byte), so any UTF-8 text round-trips "
     "losslessly \u2014 no <UNK> explosion on new characters.\n\n"
     "Vocabulary layout:\n"
     "  0-3     : <PAD>, <BOS>, <EOS>, <UNK>  (4 specials)\n"
     "  4-259   : raw bytes 0x00-0xFF         (256 base)\n"
     "  260+    : learned BPE merges\n\n"
     "Gotcha: A corpus of 'hello world' * 100 collapses into a single "
     "token after ~10 merges. BPE tests need real variation.",
     None),

    ("3 \u00b7 Sliding-window dataset \u2014 dataset.py",
     "Turns a token stream into (input_ids, targets) pairs where "
     "targets = input_ids shifted +1. This is the core causal-LM "
     "training invariant.\n\n"
     "Default stride is seq_len // 2 (50% overlap). Non-overlap gives "
     "cleaner evaluation but halves the sample count.\n\n"
     "Tested by an invariant: target[:-1] == input[1:].",
     "sliding"),

    ("4 \u00b7 The Transformer block \u2014 model.py",
     "Build bottom-up: RMSNorm \u2192 RoPE \u2192 Attention \u2192 SwiGLU \u2192 Block.\n\n"
     "Pre-Norm (normalize before the sublayer, add residual after) keeps "
     "deep stacks stable without learning-rate warmup knife-edge. Modern "
     "LLMs all use this pattern.\n\n"
     "Weight tying: self.lm_head.weight = self.token_emb.weight \u2014 not a "
     "copy, the SAME tensor. Test with 'is' not torch.equal.\n\n"
     "GPT-2 residual-projection init: scale by 1/\u221a(2\u00b7n_layers) to keep "
     "activation variance stable as depth grows.",
     "block"),

    ("5 \u00b7 Training loop \u2014 train.py",
     "The ingredients, in descending order of impact:\n\n"
     "1. bf16 mixed precision on Ampere+ \u2014 2\u00d7 speedup, basically free\n"
     "2. AdamW with weight-decay groups \u2014 decay on 2D weights only\n"
     "3. Linear warmup \u2192 cosine decay to 10% of peak LR\n"
     "4. Grad clip 1.0 \u2014 cheap insurance against loss spikes\n"
     "5. Val split is sequential 90/10 \u2014 random split leaks\n"
     "6. 'Best' checkpoint on val loss, not train loss\n\n"
     "At the end, dump loss_curve.png \u2014 students can SEE the train/val "
     "gap open up. That IS overfitting.",
     "training"),

    ("6 \u00b7 Autoregressive generation (naive)",
     "Forward the full context, sample from last-position logits, append, "
     "repeat.\n\n"
     "The sampling filter (temperature \u2192 top-k \u2192 top-p) is shared between "
     "generate(), generate_fast(), and teach.py's sampling rollout \u2014 one "
     "source of truth: _sample_from_logits().\n\n"
     "The scatter at the end of top-p looks fragile but is correct: "
     "sorted_idx is a permutation over ALL vocab positions, so every slot "
     "is written.",
     None),

    ("7 \u00b7 KV cache \u2014 generate_fast()",
     "Without a cache, every decode step recomputes Q, K, V for the entire "
     "context. With a cache, keep K, V from prior steps and only compute for "
     "the new token. Cost: O(total_len) instead of O(total_len\u00b2).\n\n"
     "Two invariants easy to miss:\n"
     "\u2022 RoPE offset: new token's angle must use start_pos = past_len, not 0\n"
     "\u2022 Causal mask: applied during prefill (T>1), skipped during decode (T=1)\n\n"
     "Correctness test: prefill(prompt) then 1-step decode must match a full "
     "forward. max |\u0394logit| < 1e-4.",
     "kv_cache"),

    ("8 \u00b7 Teaching hooks \u2014 teach.py",
     "Capture intermediate tensors via forward hooks without modifying "
     "model.py.\n\n"
     "block.attn.register_forward_pre_hook(attn_hook)\n"
     "def attn_hook(module, inp):\n"
     "    x = inp[0]\n"
     "    # Recompute Q, K, V, scores, weights from x\n"
     "    store['attn_weights'] = ...\n\n"
     "The alternative \u2014 sprinkling return_attention=True kwargs through "
     "production code \u2014 poisons model.py. Hooks keep teaching "
     "instrumentation strictly orthogonal.\n\n"
     "Re-render on-the-fly: the whole point is that changing the prompt "
     "flips the attention.",
     None),

    ("9 \u00b7 Attention visualisation \u2014 visualise.py",
     "Two things worth rendering:\n\n"
     "1. Per-head heatmap \u2014 attn_weights[layer][0, head], shape (T, T), "
     "viridis colormap, row = query, column = key.\n\n"
     "2. Attention rollout (Abnar & Zuidema 2020) \u2014 multiply attention "
     "matrices across all layers to see effective attention flow. More "
     "interpretable than any single layer.\n\n"
     "Both are in the Attention tab. Pick any prompt + layer + head to "
     "see the heatmap live.",
     None),

    ("10 \u00b7 Gradio web UI \u2014 app.py",
     "A single browser app as the webinar entrypoint. All handlers import "
     "and call existing functions \u2014 no logic re-implemented.\n\n"
     "Two Gradio patterns worth learning:\n\n"
     "1. Streaming via yield \u2014 a generator handler gives per-token "
     "typewriter output for free.\n\n"
     "2. Module-level singletons + monkeypatch for tests \u2014 load the model "
     "once at startup; tests inject a tiny model. No I/O, sub-second UI "
     "smoke tests.\n\n"
     "Version pin: gradio>=5,<6 (4.44.x has a schema-introspection bug).",
     "ui_layout"),

    ("11 \u00b7 Test suite \u2014 tests/",
     "55 tests, CPU-only, ~15s end-to-end.\n\n"
     "Highest-value tests guard invariants rather than shapes:\n\n"
     "\u2022 test_kv_cache_matches_full_pass \u2192 RoPE start_pos invariant\n"
     "\u2022 test_causal_mask_no_future_leakage \u2192 perturbing token T-1 "
     "leaves earlier logits unchanged\n"
     "\u2022 test_weight_tying \u2192 'is'-check: lm_head.weight and "
     "token_emb.weight are the SAME tensor\n"
     "\u2022 test_training_step_reduces_loss_on_tiny_overfit \u2192 end-to-end "
     "backward/optimizer moves parameters\n\n"
     "Invariant tests survive refactoring. Shape tests get rewritten "
     "every time.",
     "test_matrix"),

    ("12 \u00b7 Ship it \u2014 run.sh",
     "Commands unified through one bash dispatcher:\n\n"
     "bash run.sh setup       # venv + torch + TinyShakespeare\n"
     "bash run.sh verify      # model.py sanity check\n"
     "bash run.sh train       # full training\n"
     "bash run.sh generate --fast\n"
     "bash run.sh benchmark   # generate vs generate_fast\n"
     "bash run.sh teach       # 16 static PNGs\n"
     "bash run.sh visualise   # attention heatmaps + rollout\n"
     "bash run.sh ui          # Gradio webinar console\n"
     "bash run.sh test        # pytest suite\n\n"
     "Explicit subcommands > magic auto-detection. A new user running "
     "bash run.sh with no arg gets a helpful usage message.",
     None),
]

# ── Colors ──
BG = RGBColor(0x0F, 0x17, 0x2A)
TITLE_CLR = RGBColor(0x93, 0xC5, 0xFD)
BODY_CLR = RGBColor(0xCB, 0xD5, 0xE1)
DIM_CLR = RGBColor(0x94, 0xA3, 0xB8)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)


def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_text(slide, left, top, width, height, text,
             font_size=14, color=BODY_CLR, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        p.space_after = Pt(4)
    return box


# ── Build presentation ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
         "LLM101 \u2014 Build Your LLM From Scratch",
         font_size=36, color=TITLE_CLR, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(3.8), Inches(11), Inches(1),
         "A 12-step tour of building a ~15M-parameter GPT-style transformer "
         "in ~800 lines of PyTorch",
         font_size=18, color=BODY_CLR, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5), Inches(11), Inches(0.6),
         "RMSNorm \u00b7 RoPE \u00b7 SwiGLU \u00b7 Causal Self-Attention \u00b7 "
         "Weight Tying \u00b7 KV Cache \u00b7 BPE Tokenizer",
         font_size=14, color=DIM_CLR, align=PP_ALIGN.CENTER)

# Content slides
for title, body, img_key in STEPS:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    # Title
    add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             title, font_size=24, color=TITLE_CLR, bold=True)

    # Accent line
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.05), Inches(12.3), Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    img_path = IMG_MAP.get(img_key) if img_key else None
    has_img = img_path and os.path.exists(img_path)

    if has_img:
        # Two-column: text left, image right
        add_text(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(5.5),
                 body, font_size=13, color=BODY_CLR)
        slide.shapes.add_picture(img_path,
                                 Inches(6.3), Inches(1.3), width=Inches(6.5))
    else:
        # Full-width text
        add_text(slide, Inches(0.5), Inches(1.3), Inches(12), Inches(5.5),
                 body, font_size=15, color=BODY_CLR)

# ── Hyperparameter Effects section ──────────────────────────────
import effect_viz

effect_img_dir = tempfile.mkdtemp(prefix="llm101_effects_")
effect_paths = effect_viz.render_all(effect_img_dir)

# Section divider slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
         "Hyperparameter Effects",
         font_size=36, color=TITLE_CLR, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4), Inches(11), Inches(1),
         "How each training hyperparameter shapes the loss curve \u2014 "
         "schematic charts based on standard ML literature and observed "
         "behavior on small transformers",
         font_size=16, color=BODY_CLR, align=PP_ALIGN.CENTER)

# One slide per hyperparameter
for param in effect_viz.PARAMS:
    _, caption_md = effect_viz._PLOTS[param]
    img_path = effect_paths.get(param)

    # Strip markdown bold markers for plain text
    caption = caption_md.replace("**", "")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    # Title
    add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             f"Effect of {param}", font_size=24, color=TITLE_CLR, bold=True)

    # Accent line
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.05), Inches(12.3), Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    if img_path and os.path.exists(img_path):
        # Two-column: description left, chart right
        add_text(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(5.5),
                 caption, font_size=13, color=BODY_CLR)
        slide.shapes.add_picture(img_path,
                                 Inches(6.3), Inches(1.3), width=Inches(6.5))
    else:
        add_text(slide, Inches(0.5), Inches(1.3), Inches(12), Inches(5.5),
                 caption, font_size=15, color=BODY_CLR)

# ── Save ──
out_path = os.path.join(os.path.dirname(__file__), "LLM101_Build_Steps.pptx")
prs.save(out_path)
n_effect = len(effect_viz.PARAMS)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)} (1 title + 12 steps + 1 section + {n_effect} effects)")
