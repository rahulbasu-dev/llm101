"""LLM101 - Console (Gradio UI).

A single browser app with seven tabs:
  1. Generate      — live token streaming; toggle KV cache on/off to see speed
  2. Train Reports — render the 16 step-by-step slides for a chosen prompt
  3. Attention     — per-head heatmap + rollout for the current prompt
  4. Benchmark     — side-by-side generate() vs generate_fast() tok/s
  5. Train         — interactive training with LR / dropout / warmup sliders
  6. Effects       — how each hyperparameter shapes the loss curve
  7. Build Steps   — structural diagrams (no trained model needed)

All tabs reuse the existing teach.py / visualise.py / build_viz.py /
effect_viz.py plumbing — no logic is duplicated. The model is loaded once
at startup.

Run:
    python app.py
    python app.py --share           # phone-home for a public URL
    python app.py --port 7861
    bash run.sh ui
"""

from __future__ import annotations
import argparse
import json
import math
import os
import socket
import tempfile
import time
import warnings

# Silence Gradio 6.0 deprecation warnings that Gradio 5.x emits for APIs that
# still work fine in 5.x (theme= on Blocks, show_api= on launch). Targeted
# pattern avoids swallowing unrelated warnings.
warnings.filterwarnings(
    "ignore",
    message=r".*will be removed in Gradio 6\.0.*",
    category=DeprecationWarning,
)

import gradio as gr
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import torch
import torch.nn.functional as F

from config import NanoLLMConfig, require_cuda
from tokenizer import BPETokenizer
from model import NanoLLM, _sample_from_logits
from train import train_iter, finetune_iter, save_loss_curve
from forward_viz import (
    ForwardCapture, token_labels, _get_plt,
    slide_01_tokenization, slide_02_embeddings, slide_03_qkv,
    slide_04_scores_raw, slide_05_scores_masked, slide_06_attn_weights,
    slide_07_value_sum, slide_08_all_heads, slide_09_ffn_delta,
    slide_10_logits_topk, slide_11_sampling_rollout,
    slide_12_positional_rope, slide_13_scaling_rationale,
    slide_14_temperature_effect, slide_15_greedy_vs_sample,
    slide_16_param_breakdown,
)
from visualise import AttentionCapture, compute_attention_rollout, plot_attention_heatmap
from visualize_anim import collect_viz_data
import build_viz
import effect_viz
import notebook_viz as nv


# ═══════════════════════════════════════════════════════════════
# Model loading (singleton — load once, share across handlers)
# ═══════════════════════════════════════════════════════════════

_MODEL: NanoLLM | None = None
_TOKENIZER: BPETokenizer | None = None
_CONFIG: NanoLLMConfig | None = None
_STATUS: str = ""  # Banner text shown in the header


def _device_badge(device) -> str:
    """Return an HTML-coloured device label for the status banner."""
    label = str(device)
    if label.startswith("cuda"):
        return f'<span style="color:#4ade80;font-weight:bold">{label}</span>'
    return f'<span style="color:#f87171;font-weight:bold">{label} ⚠ CPU — training will be slow</span>'


def _load_model(checkpoint_path: str = "checkpoints/best.pt") -> None:
    """Load model + tokenizer once. Falls back to random weights if no checkpoint."""
    global _MODEL, _TOKENIZER, _CONFIG, _STATUS

    device = require_cuda()  # hard-require CUDA (set NANOLLM_ALLOW_CPU=1 to bypass)

    if os.path.exists(checkpoint_path):
        ckpt = torch.load(checkpoint_path, map_location=device, weights_only=False)
        config = ckpt["config"]
        tokenizer = BPETokenizer(target_vocab_size=config.target_vocab_size)
        tokenizer.load(config.tokenizer_path)
        config.vocab_size = tokenizer.vocab_size
        model = NanoLLM(config).to(device)
        # Checkpoints saved from torch.compile'd models have keys prefixed
        # with "_orig_mod." — strip that prefix so they load into a plain model.
        sd = ckpt["model_state_dict"]
        if any(k.startswith("_orig_mod.") for k in sd):
            sd = {k.removeprefix("_orig_mod."): v for k, v in sd.items()}
        model.load_state_dict(sd)
        model.eval()

        epoch = ckpt.get("epoch", "?")
        val_loss = ckpt.get("val_loss", ckpt.get("loss", "?"))
        _STATUS = (
            f"**Loaded:** `{checkpoint_path}` (epoch {epoch}, val_loss={val_loss}) "
            f"· **Device:** {_device_badge(device)} · **Vocab:** {config.vocab_size} "
            f"· **Params:** {sum(p.numel() for p in model.parameters())/1e6:.2f}M"
        )
    else:
        # Fall back to a random-initialised tiny model so the UI is explorable
        # even before training finishes.
        config = NanoLLMConfig(vocab_size=260)  # Just the byte vocab
        tokenizer = BPETokenizer(target_vocab_size=260)
        model = NanoLLM(config).to(device).eval()
        _STATUS = (
            f"⚠️ **No checkpoint at `{checkpoint_path}`** — using random weights. "
            f"Generation will be garbage. Run `bash run.sh train` first for real output. "
            f"· **Device:** {_device_badge(device)}"
        )

    _MODEL = model
    _TOKENIZER = tokenizer
    _CONFIG = config
    _restore_train_state()


def _require_loaded():
    """Assert the model has been loaded. All handlers depend on this."""
    if _MODEL is None:
        raise RuntimeError("Model not loaded. Call _load_model() first.")
    return _MODEL, _TOKENIZER, _CONFIG


# ═══════════════════════════════════════════════════════════════
# Tab 1: Generate (with token-by-token streaming)
# ═══════════════════════════════════════════════════════════════

def generate_stream(prompt, max_new_tokens, temperature, top_k, top_p, use_cache):
    """Yield the accumulated decoded text after each new token.

    Gradio will update the Textbox once per yield, producing the classic
    typewriter effect that makes autoregressive generation tangible to
    students: they SEE each token appearing one at a time.
    """
    model, tokenizer, config = _require_loaded()
    if not prompt:
        yield "(empty prompt)"
        return

    ids = tokenizer.encode(prompt, add_special=False)
    if not ids:
        yield "(prompt encoded to zero tokens)"
        return

    # Respect context window
    if len(ids) >= config.max_seq_len:
        ids = ids[-(config.max_seq_len - 1):]

    idx = torch.tensor([ids], device=config.device)
    out_ids = list(ids)
    yield tokenizer.decode(out_ids)  # initial state

    t0 = time.time()

    with torch.no_grad():
        if use_cache:
            # Prefill + incremental decode with KV cache
            logits, cache = model(idx)
            for step in range(int(max_new_tokens)):
                if cache[0][0].size(2) >= config.max_seq_len:
                    break
                next_tok = _sample_from_logits(logits, temperature, int(top_k), top_p)
                out_ids.append(int(next_tok.item()))
                yield tokenizer.decode(out_ids)
                logits, cache = model(next_tok, past_kv=cache)
        else:
            # Reference path: re-process the whole context every step
            cur = idx
            for step in range(int(max_new_tokens)):
                crop = cur[:, -config.max_seq_len:]
                logits, _ = model(crop)
                next_tok = _sample_from_logits(logits, temperature, int(top_k), top_p)
                cur = torch.cat([cur, next_tok], dim=1)
                out_ids.append(int(next_tok.item()))
                yield tokenizer.decode(out_ids)

    # Final frame with timing (use the known step count, not the loop var,
    # which may not be in scope if max_new_tokens was 0).
    elapsed = time.time() - t0
    n_new = len(out_ids) - len(ids)
    if n_new > 0:
        tps = n_new / max(elapsed, 1e-6)
        final_text = tokenizer.decode(out_ids) + \
            f"\n\n— {n_new} tokens in {elapsed:.2f}s ({tps:.1f} tok/s, " + \
            ("with cache" if use_cache else "no cache") + ")"
        yield final_text


# ═══════════════════════════════════════════════════════════════
# Tab 4: Train Reports (render 16 slides for any prompt)
# ═══════════════════════════════════════════════════════════════

def render_teach(prompt, layer, head, query_pos, temperature, top_k, top_p):
    """Render the 16-slide walkthrough for a user-chosen prompt + layer + head."""
    model, tokenizer, config = _require_loaded()
    plt_ = _get_plt()

    if not prompt:
        return [], "Please enter a prompt."

    token_ids = tokenizer.encode(prompt, add_special=False)
    if len(token_ids) < 2:
        return [], f"Prompt tokenizes to only {len(token_ids)} token(s). Need ≥ 2."
    if len(token_ids) > config.max_seq_len:
        token_ids = token_ids[:config.max_seq_len]
        prompt = tokenizer.decode(token_ids)

    # Resolve indices
    query_pos = int(query_pos)
    if query_pos < 0:
        query_pos = len(token_ids) - 1
    query_pos = min(max(0, query_pos), len(token_ids) - 1)
    layer = min(max(0, int(layer)), config.n_layers - 1)
    head = min(max(0, int(head)), config.n_heads - 1)

    labels = token_labels(tokenizer, token_ids, max_len=8)
    idx = torch.tensor([token_ids], device=config.device)

    # Run forward with hooks
    capture = ForwardCapture(model, target_layer=layer)
    with torch.no_grad():
        model(idx)
    capture.remove()

    # Write PNGs to a fresh temp dir (avoid caching old renders)
    outdir = tempfile.mkdtemp(prefix="nanollm_ui_slides_")
    p = lambda n: os.path.join(outdir, n)

    slide_01_tokenization(plt_, tokenizer, prompt, token_ids, p("01_tokenization.png"))
    slide_02_embeddings(plt_, capture.store, labels, p("02_embeddings.png"))
    slide_03_qkv(plt_, capture.store, labels, head, p("03_qkv.png"))
    slide_04_scores_raw(plt_, capture.store, labels, head, p("04_scores_raw.png"))
    slide_05_scores_masked(plt_, capture.store, labels, head, p("05_scores_masked.png"))
    slide_06_attn_weights(plt_, capture.store, labels, head, p("06_attn_weights.png"))
    slide_07_value_sum(plt_, capture.store, labels, head, query_pos, p("07_value_sum.png"))
    slide_08_all_heads(plt_, capture.store, labels, p("08_all_heads.png"))
    slide_09_ffn_delta(plt_, capture.store, labels, p("09_ffn_delta.png"))
    slide_10_logits_topk(plt_, model, tokenizer, idx, config.device,
                         temperature, int(top_k), top_p, p("10_logits_topk.png"))
    slide_11_sampling_rollout(plt_, model, tokenizer, idx, config.device,
                              temperature, int(top_k), top_p, p("11_sampling_rollout.png"))
    slide_12_positional_rope(plt_, model, p("12_positional_rope.png"))
    slide_13_scaling_rationale(plt_, p("13_scaling_rationale.png"))
    slide_14_temperature_effect(plt_, model, tokenizer, idx, p("14_temperature_effect.png"))
    slide_15_greedy_vs_sample(plt_, model, tokenizer, idx, config.device, p("15_greedy_vs_sample.png"))
    slide_16_param_breakdown(plt_, model, p("16_param_breakdown.png"))

    # Build gallery entries: (image_path, caption)
    captions = [
        "01 · Tokenization — text → bytes → BPE → token IDs",
        "02 · Embeddings — token IDs → dense vectors",
        f"03 · Q, K, V — projections for layer head {head}",
        f"04 · Attention scores (raw, before mask)",
        "05 · Causal mask applied (upper triangle → −∞)",
        "06 · Softmax → attention weights (rows sum to 1)",
        f"07 · Weighted value sum for query '{labels[query_pos]}'",
        "08 · All heads side-by-side (this layer)",
        "09 · FFN: before vs delta vs after",
        "10 · Next-token distribution (top 20)",
        "11 · Sampling rollout — 10 decode steps",
        "12 · RoPE positional encoding",
        "13 · Why /√d_head (scaling rationale)",
        "14 · Temperature effect on distribution",
        "15 · Greedy vs sampling output",
        "16 · Parameter breakdown",
    ]
    gallery = list(zip(sorted(os.listdir(outdir)), captions))
    gallery = [(os.path.join(outdir, f), c) for f, c in gallery]

    status = (
        f"Rendered 16 slides for prompt '{prompt[:40]}{'…' if len(prompt)>40 else ''}' "
        f"({len(token_ids)} tokens) · layer={layer} · head={head} · "
        f"query_pos={query_pos} ('{labels[query_pos]}')"
    )
    return gallery, status, outdir


def export_slides_pptx(slide_dir: str, params_label: str) -> str | None:
    """Export rendered slide PNGs from a directory to a PPTX file."""
    try:
        from pptx import Presentation as Pptx
        from pptx.util import Inches as In, Pt as PtU
        from pptx.dml.color import RGBColor as RGB
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return None

    if not slide_dir or not os.path.isdir(slide_dir):
        return None

    prs = Pptx()
    prs.slide_width = In(13.333)
    prs.slide_height = In(7.5)

    BG = RGB(0x0F, 0x17, 0x2A)
    TITLE_CLR = RGB(0x93, 0xC5, 0xFD)
    DIM_CLR = RGB(0x94, 0xA3, 0xB8)

    # Title slide
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = BG
    tb = sl.shapes.add_textbox(In(0.5), In(2.5), In(12), In(1))
    p = tb.text_frame.paragraphs[0]
    p.text = "Train Reports — 16-slide forward-pass walkthrough"
    p.font.size = PtU(28)
    p.font.color.rgb = TITLE_CLR
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tb2 = sl.shapes.add_textbox(In(0.5), In(3.8), In(12), In(0.6))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = params_label
    p2.font.size = PtU(16)
    p2.font.color.rgb = DIM_CLR
    p2.alignment = PP_ALIGN.CENTER

    # Slide captions
    captions = [
        "01 · Tokenization", "02 · Embeddings", "03 · Q, K, V projections",
        "04 · Attention scores (raw)", "05 · Causal mask applied",
        "06 · Softmax attention weights", "07 · Weighted value sum",
        "08 · All heads side-by-side", "09 · FFN before/delta/after",
        "10 · Next-token distribution", "11 · Sampling rollout",
        "12 · RoPE positional encoding", "13 · Scaling rationale",
        "14 · Temperature effect", "15 · Greedy vs sampling",
        "16 · Parameter breakdown",
    ]

    pngs = sorted(f for f in os.listdir(slide_dir) if f.endswith(".png"))
    for i, png in enumerate(pngs):
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        sl.background.fill.solid()
        sl.background.fill.fore_color.rgb = BG
        # Caption
        cap = captions[i] if i < len(captions) else png
        tb = sl.shapes.add_textbox(In(0.3), In(0.2), In(12.5), In(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = cap
        p.font.size = PtU(20)
        p.font.color.rgb = TITLE_CLR
        p.font.bold = True
        # Image centered
        sl.shapes.add_picture(
            os.path.join(slide_dir, png),
            In(0.5), In(0.9), width=In(12.3),
        )

    out_path = os.path.join(slide_dir, "Train_Reports.pptx")
    prs.save(out_path)
    return out_path


# State for pin-and-compare
_PINNED_GALLERY = []
_PINNED_LABEL = ""
_LAST_SLIDE_DIR = None


# ═══════════════════════════════════════════════════════════════
# Tab 3: Attention (one head + rollout)
# ═══════════════════════════════════════════════════════════════

def render_attention(prompt, layer, head):
    """Render a single head's heatmap + the attention-rollout heatmap."""
    model, tokenizer, config = _require_loaded()

    if not prompt:
        return None, None, "Please enter a prompt."

    token_ids = tokenizer.encode(prompt, add_special=False)
    if len(token_ids) < 2:
        return None, None, f"Prompt tokenizes to only {len(token_ids)} token(s). Need ≥ 2."
    if len(token_ids) > config.max_seq_len:
        token_ids = token_ids[:config.max_seq_len]

    labels = [tokenizer.decode_token(t).replace("\n", "\\n")[:10] for t in token_ids]
    idx = torch.tensor([token_ids], device=config.device)

    cap = AttentionCapture(model)
    with torch.no_grad():
        model(idx)
    cap.remove_hooks()

    n_layers = len(cap.attention_maps)
    n_heads = cap.attention_maps[0].size(1)
    layer = min(max(0, int(layer)), n_layers - 1)
    head = min(max(0, int(head)), n_heads - 1)

    outdir = tempfile.mkdtemp(prefix="nanollm_ui_attn_")
    head_path = os.path.join(outdir, "head.png")
    rollout_path = os.path.join(outdir, "rollout.png")

    # Single head heatmap
    weights = cap.attention_maps[layer][0, head].numpy()
    plot_attention_heatmap(
        weights, labels, f"Layer {layer}, Head {head}", head_path
    )

    # Rollout across all layers
    rollout = compute_attention_rollout(cap.attention_maps)
    plot_attention_heatmap(
        rollout[0].numpy(), labels,
        f"Attention rollout across {n_layers} layers",
        rollout_path,
    )

    info = (
        f"{len(token_ids)} tokens · layer {layer}/{n_layers-1} · "
        f"head {head}/{n_heads-1}"
    )
    return head_path, rollout_path, info


# ═══════════════════════════════════════════════════════════════
# Tab 4: Benchmark (generate vs generate_fast)
# ═══════════════════════════════════════════════════════════════

def run_bench(prompt, n_tokens):
    """Time both generation paths and return a bar chart + summary text."""
    model, tokenizer, config = _require_loaded()
    n_tokens = int(n_tokens)

    prompt = prompt or "The "
    ids = tokenizer.encode(prompt, add_special=False) or [0]
    idx = torch.tensor([ids], device=config.device)

    # Warm up once so kernel-compile cost isn't charged to the first path
    with torch.no_grad():
        model.generate(idx, max_new_tokens=4, temperature=0.8, top_k=40, top_p=0.9)
        model.generate_fast(idx, max_new_tokens=4, temperature=0.8, top_k=40, top_p=0.9)
    if config.device.type == "cuda":
        torch.cuda.synchronize()

    # No-cache path
    t0 = time.time()
    with torch.no_grad():
        model.generate(idx, max_new_tokens=n_tokens)
    if config.device.type == "cuda":
        torch.cuda.synchronize()
    t_slow = time.time() - t0

    # KV-cache path
    t0 = time.time()
    with torch.no_grad():
        model.generate_fast(idx, max_new_tokens=n_tokens)
    if config.device.type == "cuda":
        torch.cuda.synchronize()
    t_fast = time.time() - t0

    slow_tps = n_tokens / t_slow
    fast_tps = n_tokens / t_fast
    speedup = t_slow / t_fast if t_fast > 0 else float("inf")

    # Bar chart
    fig, ax = plt.subplots(figsize=(8, 4.5))
    bars = ax.bar(
        ["generate()\n(no cache)", "generate_fast()\n(KV cache)"],
        [slow_tps, fast_tps],
        color=["#9ecae1", "#08519c"],
    )
    ax.set_ylabel("tokens / second", fontsize=11)
    ax.set_title(
        f"Generation throughput — {n_tokens} tokens · speedup {speedup:.2f}×",
        fontsize=12, fontweight="bold",
    )
    for bar, val in zip(bars, [slow_tps, fast_tps]):
        ax.text(bar.get_x() + bar.get_width() / 2, val * 1.02,
                f"{val:.1f} tok/s", ha="center", va="bottom",
                fontsize=11, fontweight="bold")
    ax.grid(True, axis="y", alpha=0.3)
    try:
        fig.tight_layout()
    except (ValueError, Exception):
        pass

    outpath = os.path.join(tempfile.mkdtemp(prefix="nanollm_ui_bench_"), "bench.png")
    fig.savefig(outpath, dpi=130)
    plt.close(fig)

    summary = (
        f"generate()        : {t_slow*1000:8.1f} ms   ·   {slow_tps:6.1f} tok/s\n"
        f"generate_fast()   : {t_fast*1000:8.1f} ms   ·   {fast_tps:6.1f} tok/s\n"
        f"Speedup           : {speedup:.2f}x"
    )
    return outpath, summary


# ═══════════════════════════════════════════════════════════════
# Tab 5: Train (streams training progress into the UI)
# ═══════════════════════════════════════════════════════════════

# ── Background-training shared state ──────────────────────────
# Training runs in a daemon thread so it survives tab switches and
# browser refreshes. The UI polls _TRAIN_STATE every 2 s via gr.Timer.
import threading as _threading

_TRAIN_STATE: dict = {
    "running": False,
    "log_lines": [],
    "plot_path": None,
    "status": "idle",
    "samples": "",
    "train_history": [],
    "epoch_history": [],
}
_TRAIN_LOCK = _threading.Lock()
_TRAIN_STOP = _threading.Event()

_TRAIN_STATE_PATH = os.path.join("checkpoints", "train_state.json")


def _save_train_state() -> None:
    """Persist log, history, samples, and status to disk (called after each epoch)."""
    try:
        os.makedirs("checkpoints", exist_ok=True)
        with _TRAIN_LOCK:
            payload = {
                "status":        _TRAIN_STATE["status"],
                "samples":       _TRAIN_STATE["samples"],
                "log_lines":     _TRAIN_STATE["log_lines"][-500:],  # keep last 500 lines
                "train_history": _TRAIN_STATE["train_history"],
                "epoch_history": _TRAIN_STATE["epoch_history"],
            }
        with open(_TRAIN_STATE_PATH, "w", encoding="utf-8") as f:
            json.dump(payload, f, indent=2)
    except Exception:
        pass  # persistence is best-effort; never crash training


def _restore_train_state() -> None:
    """Load persisted training state on app startup if a saved run exists."""
    if not os.path.exists(_TRAIN_STATE_PATH):
        return
    try:
        with open(_TRAIN_STATE_PATH, "r", encoding="utf-8") as f:
            payload = json.load(f)
        curve = os.path.join("checkpoints", "loss_curve.png")
        with _TRAIN_LOCK:
            _TRAIN_STATE.update({
                "running":       False,
                "status":        payload.get("status", "done"),
                "samples":       payload.get("samples", ""),
                "log_lines":     payload.get("log_lines", []),
                "train_history": [tuple(x) for x in payload.get("train_history", [])],
                "epoch_history": [tuple(x) for x in payload.get("epoch_history", [])],
                "plot_path":     curve if os.path.exists(curve) else None,
            })
    except Exception:
        pass


def _build_live_curve(history: list[tuple[int, float]],
                      epochs: list[tuple[int, float, float]]) -> str | None:
    """Render an in-progress loss curve PNG. Returns the path or None."""
    if not history:
        return None
    outdir = tempfile.mkdtemp(prefix="nanollm_train_curve_")
    path = os.path.join(outdir, "curve.png")
    save_loss_curve(history, epochs, path)
    return path


def _async_curve(history: list, epochs: list) -> None:
    """Render loss curve in a daemon thread so training isn't blocked."""
    path = _build_live_curve(history, epochs)
    if path:
        with _TRAIN_LOCK:
            _TRAIN_STATE["plot_path"] = path


def _refresh_curve(history: list, epochs: list) -> None:
    _threading.Thread(
        target=_async_curve, args=(list(history), list(epochs)), daemon=True
    ).start()


def _train_bg(cfg: "NanoLLMConfig",
              checkpoint_path: str = "checkpoints/best.pt") -> None:
    """Background thread: runs train_iter and writes into _TRAIN_STATE."""
    global _MODEL, _TOKENIZER, _CONFIG
    step_counter = 0
    _t0 = time.time()

    def _elapsed() -> str:
        s = int(time.time() - _t0)
        return f"{s // 3600:02d}:{(s % 3600) // 60:02d}:{s % 60:02d}"

    def _log(msg: str) -> None:
        with _TRAIN_LOCK:
            _TRAIN_STATE["log_lines"].append(msg)

    try:
        for evt in train_iter(cfg, stop_event=_TRAIN_STOP):
            t = evt["type"]

            if t == "log":
                _log(evt["msg"])

            elif t == "step":
                step_counter += 1
                with _TRAIN_LOCK:
                    _TRAIN_STATE["train_history"].append(
                        (evt["global_step"], evt["loss"]))
                    _hist_snap = list(_TRAIN_STATE["train_history"])
                    _ep_snap   = list(_TRAIN_STATE["epoch_history"])
                if step_counter % 10 == 0:
                    _refresh_curve(_hist_snap, _ep_snap)
                if evt["batch_idx"] % cfg.log_interval == 0:
                    ppl = (math.exp(min(evt["loss"], 20))
                           if evt["loss"] < 20 else float("inf"))
                    st = evt.get("step_time", 0)
                    _log(f"  [{_elapsed()}] Epoch {evt['epoch']}/{cfg.max_epochs} | "
                         f"Step {evt['batch_idx']}/{evt['total_batches']} | "
                         f"Loss {evt['loss']:.4f} | PPL {ppl:.1f} | "
                         f"LR {evt['lr']:.2e} | {evt['tps']:,.0f} tok/s | {st:.1f}s/step")

            elif t == "epoch":
                with _TRAIN_LOCK:
                    _TRAIN_STATE["epoch_history"].append(
                        (evt["epoch"], evt["train_loss"], evt["val_loss"]))
                    _hist_snap = list(_TRAIN_STATE["train_history"])
                    _ep_snap   = list(_TRAIN_STATE["epoch_history"])
                mins, secs = divmod(int(evt["elapsed"]), 60)
                _log("")
                _log(f"  Epoch {evt['epoch']}/{evt['max_epochs']} done  |  "
                     f"train={evt['train_loss']:.4f}  val={evt['val_loss']:.4f}  "
                     f"(ppl {evt['train_ppl']:.1f} / {evt['val_ppl']:.1f})  "
                     f"epoch time={mins}m{secs:02d}s")
                if evt["samples"]:
                    preview = evt["samples"][0].replace("\n", " / ")[:200]
                    new_entry = f"Epoch {evt['epoch']} sample:\n  \"{preview}\""
                    with _TRAIN_LOCK:
                        existing = _TRAIN_STATE["samples"]
                        _TRAIN_STATE["samples"] = (
                            new_entry if not existing
                            else new_entry + "\n\n" + existing)
                _refresh_curve(_hist_snap, _ep_snap)
                _threading.Thread(target=_save_train_state, daemon=True).start()

            elif t == "best":
                _log(f"  * New best val loss {evt['val_loss']:.4f} -> {evt['path']}")

            elif t == "done":
                _log("")
                _log("=" * 50)
                _log(f"Training complete.  Best val loss: {evt['best_val_loss']:.4f}")
                _log(f"Loss curve: {evt['curve_path']}")
                _log("=" * 50)
                with _TRAIN_LOCK:
                    if evt["curve_path"]:
                        _TRAIN_STATE["plot_path"] = evt["curve_path"]

            elif t == "stopped":
                _log("Training stopped by user.")
                with _TRAIN_LOCK:
                    _TRAIN_STATE["status"] = "stopped"
                _save_train_state()
                return

            elif t == "error":
                _log(f"ERROR: {evt['msg']}")
                with _TRAIN_LOCK:
                    _TRAIN_STATE["status"] = "error"
                    _TRAIN_STATE["running"] = False
                return

        try:
            _load_model(checkpoint_path)
            _log(f"Model reloaded.  Status: {_STATUS}")
        except Exception as exc:
            _log(f"  (note: couldn't auto-reload model: {exc})")

        with _TRAIN_LOCK:
            _TRAIN_STATE["status"] = "done"
        _save_train_state()

    except Exception as exc:
        _log(f"FATAL: {exc}")
        with _TRAIN_LOCK:
            _TRAIN_STATE["status"] = "error"
    finally:
        with _TRAIN_LOCK:
            _TRAIN_STATE["running"] = False


def train_start(max_epochs_override: int, batch_size_override: int,
                learning_rate_override: float, dropout_override: float,
                warmup_steps_override: int,
                checkpoint_path: str = "checkpoints/best.pt") -> str:
    """Launch training in a background thread; return immediately."""
    with _TRAIN_LOCK:
        if _TRAIN_STATE["running"]:
            return "already running"
        _TRAIN_STOP.clear()
        _TRAIN_STATE.update({
            "running": True,
            "log_lines": [],
            "plot_path": None,
            "status": "running",
            "samples": "",
            "train_history": [],
            "epoch_history": [],
        })

    cfg = NanoLLMConfig()
    cfg.max_epochs = int(max_epochs_override)
    cfg.batch_size = int(batch_size_override)
    cfg.learning_rate = float(learning_rate_override)
    cfg.dropout = float(dropout_override)
    if int(warmup_steps_override) > 0:
        cfg.warmup_steps = int(warmup_steps_override)
        setattr(cfg, "_warmup_explicit", True)

    _threading.Thread(
        target=_train_bg, args=(cfg, checkpoint_path), daemon=True
    ).start()
    return "running"


def poll_train() -> tuple:
    """Called by gr.Timer every 2 s — returns latest training state."""
    with _TRAIN_LOCK:
        tail = _TRAIN_STATE["log_lines"][-300:]
        log = "\n".join(tail)
        plot = _TRAIN_STATE["plot_path"]
        status = _TRAIN_STATE["status"]
        samples = _TRAIN_STATE["samples"] or "(samples appear here after each epoch completes)"
    return log, plot, status, samples


def train_stop() -> str:
    """Signal the background training thread to stop after the current batch."""
    _TRAIN_STOP.set()
    return "stopping…"


def train_restart(max_epochs_override: int, batch_size_override: int,
                  learning_rate_override: float, dropout_override: float,
                  warmup_steps_override: int,
                  checkpoint_path: str = "checkpoints/best.pt") -> str:
    """Stop any running training then immediately start a fresh run."""
    _TRAIN_STOP.set()
    # Brief spin-wait (max 3 s) for the thread to notice and exit
    for _ in range(30):
        with _TRAIN_LOCK:
            if not _TRAIN_STATE["running"]:
                break
        _threading.Event().wait(0.1)
    return train_start(max_epochs_override, batch_size_override,
                       learning_rate_override, dropout_override,
                       warmup_steps_override, checkpoint_path)


# ═══════════════════════════════════════════════════════════════
# Fine-tune tab: adapt the pre-trained model on custom text
# ═══════════════════════════════════════════════════════════════

_FT_STATE: dict = {
    "running": False,
    "log_lines": [],
    "plot_path": None,
    "status": "idle",
    "samples": "",
    "train_history": [],
    "epoch_history": [],
}
_FT_LOCK = _threading.Lock()
_FT_STOP = _threading.Event()


def _ft_bg(custom_text: str, checkpoint_path: str,
           learning_rate: float, max_epochs: int,
           batch_size: int) -> None:
    """Background thread: runs finetune_iter and writes into _FT_STATE."""
    global _MODEL, _TOKENIZER, _CONFIG
    step_counter = 0
    _t0 = time.time()

    def _elapsed() -> str:
        s = int(time.time() - _t0)
        return f"{s // 3600:02d}:{(s % 3600) // 60:02d}:{s % 60:02d}"

    def _log(msg: str) -> None:
        with _FT_LOCK:
            _FT_STATE["log_lines"].append(msg)

    try:
        ft_cfg = NanoLLMConfig()
        ft_cfg.learning_rate = learning_rate
        ft_cfg.max_epochs = max_epochs
        ft_cfg.batch_size = batch_size

        for evt in finetune_iter(
            custom_text,
            checkpoint_path=checkpoint_path,
            config=ft_cfg,
            stop_event=_FT_STOP,
        ):
            t = evt["type"]

            if t == "log":
                _log(evt["msg"])

            elif t == "step":
                step_counter += 1
                with _FT_LOCK:
                    _FT_STATE["train_history"].append(
                        (evt["global_step"], evt["loss"]))
                    _hist_snap = list(_FT_STATE["train_history"])
                    _ep_snap   = list(_FT_STATE["epoch_history"])
                if step_counter % 10 == 0:
                    path = _build_live_curve(_hist_snap, _ep_snap)
                    if path:
                        with _FT_LOCK:
                            _FT_STATE["plot_path"] = path
                if evt["batch_idx"] % 5 == 0:
                    ppl = (math.exp(min(evt["loss"], 20))
                           if evt["loss"] < 20 else float("inf"))
                    _log(f"  [{_elapsed()}] Epoch {evt['epoch']}/{max_epochs} | "
                         f"Step {evt['batch_idx']}/{evt['total_batches']} | "
                         f"Loss {evt['loss']:.4f} | PPL {ppl:.1f} | "
                         f"LR {evt['lr']:.2e} | {evt['tps']:,.0f} tok/s")

            elif t == "epoch":
                with _FT_LOCK:
                    _FT_STATE["epoch_history"].append(
                        (evt["epoch"], evt["train_loss"], evt["val_loss"]))
                    _hist_snap = list(_FT_STATE["train_history"])
                    _ep_snap   = list(_FT_STATE["epoch_history"])
                mins, secs = divmod(int(evt["elapsed"]), 60)
                _log("")
                _log(f"  Epoch {evt['epoch']}/{evt['max_epochs']} done  |  "
                     f"train={evt['train_loss']:.4f}  val={evt['val_loss']:.4f}  "
                     f"(ppl {evt['train_ppl']:.1f} / {evt['val_ppl']:.1f})  "
                     f"epoch time={mins}m{secs:02d}s")
                if evt["samples"]:
                    preview = evt["samples"][0].replace("\n", " / ")[:200]
                    new_entry = f"Epoch {evt['epoch']} sample:\n  \"{preview}\""
                    with _FT_LOCK:
                        existing = _FT_STATE["samples"]
                        _FT_STATE["samples"] = (
                            new_entry if not existing
                            else new_entry + "\n\n" + existing)
                path = _build_live_curve(_hist_snap, _ep_snap)
                if path:
                    with _FT_LOCK:
                        _FT_STATE["plot_path"] = path

            elif t == "best":
                _log(f"  * New best val loss {evt['val_loss']:.4f} -> {evt['path']}")

            elif t == "done":
                _log("")
                _log("=" * 50)
                _log(f"Fine-tuning complete.  Best val loss: {evt['best_val_loss']:.4f}")
                _log(f"Saved: {evt.get('curve_path', 'checkpoints/finetuned.pt')}")
                _log("=" * 50)
                with _FT_LOCK:
                    if evt.get("curve_path"):
                        _FT_STATE["plot_path"] = evt["curve_path"]

            elif t == "stopped":
                _log("Fine-tuning stopped by user.")
                with _FT_LOCK:
                    _FT_STATE["status"] = "stopped"
                return

            elif t == "error":
                _log(f"ERROR: {evt['msg']}")
                with _FT_LOCK:
                    _FT_STATE["status"] = "error"
                    _FT_STATE["running"] = False
                return

        # Reload the fine-tuned model so Generate tab uses updated weights
        try:
            ft_path = "checkpoints/finetuned.pt"
            if os.path.exists(ft_path):
                _load_model(ft_path)
                _log(f"Fine-tuned model reloaded from {ft_path}")
        except Exception as exc:
            _log(f"  (note: couldn't auto-reload fine-tuned model: {exc})")

        with _FT_LOCK:
            _FT_STATE["status"] = "done"

    except Exception as exc:
        _log(f"FATAL: {exc}")
        with _FT_LOCK:
            _FT_STATE["status"] = "error"
    finally:
        with _FT_LOCK:
            _FT_STATE["running"] = False


def ft_start(custom_text: str, checkpoint_path: str,
             learning_rate: float, max_epochs: int,
             batch_size: int) -> str:
    """Launch fine-tuning in a background thread; return immediately."""
    if not custom_text or len(custom_text.strip()) < 200:
        return "error: paste at least 200 characters of text to fine-tune on"
    with _FT_LOCK:
        if _FT_STATE["running"]:
            return "already running"
        _FT_STOP.clear()
        _FT_STATE.update({
            "running": True,
            "log_lines": [],
            "plot_path": None,
            "status": "running",
            "samples": "",
            "train_history": [],
            "epoch_history": [],
        })
    _threading.Thread(
        target=_ft_bg,
        args=(custom_text, checkpoint_path, learning_rate, int(max_epochs), int(batch_size)),
        daemon=True,
    ).start()
    return "running"


def ft_stop() -> str:
    """Signal the fine-tuning thread to stop after the current batch."""
    _FT_STOP.set()
    return "stopping…"


def poll_ft() -> tuple:
    """Called by gr.Timer every 2 s — returns latest fine-tune state."""
    with _FT_LOCK:
        tail = _FT_STATE["log_lines"][-300:]
        log = "\n".join(tail)
        plot = _FT_STATE["plot_path"]
        status = _FT_STATE["status"]
        samples = _FT_STATE["samples"] or "(samples appear here after each epoch completes)"
    return log, plot, status, samples


# ═══════════════════════════════════════════════════════════════
# Tab 6: Effects (how each hyperparameter shapes training)
# ═══════════════════════════════════════════════════════════════

def render_effect_panel(param: str) -> tuple[str, str]:
    """Return (markdown_caption, image_path) for one effect selection."""
    outdir = globals().setdefault(
        "_EFFECT_DIR",
        tempfile.mkdtemp(prefix="nanollm_effects_"),
    )
    img, caption = effect_viz.render(param, outdir)
    return caption, img


# ═══════════════════════════════════════════════════════════════
# Tab 7: Build Steps (interactive tutorial)
# ═══════════════════════════════════════════════════════════════

# ═══════════════════════════════════════════════════════════════
# Tab 7: Visualize (animated forward pass)
# ═══════════════════════════════════════════════════════════════

_VIZ_TEMPLATE = None  # Lazy-loaded


def _get_viz_template() -> str:
    """Load the HTML template once."""
    global _VIZ_TEMPLATE
    if _VIZ_TEMPLATE is None:
        template_path = os.path.join(os.path.dirname(__file__), "templates", "visualize.html")
        with open(template_path, "r", encoding="utf-8") as f:
            _VIZ_TEMPLATE = f.read()
    return _VIZ_TEMPLATE


def render_visualization(prompt: str) -> str:
    """Collect tensors and return the animated HTML visualization.

    Gradio's gr.HTML does NOT execute <script> tags (static HTML only).
    We wrap the visualization in an <iframe srcdoc="..."> which creates
    a sandboxed document that does execute scripts.
    """
    import json
    import html as html_lib
    model, tokenizer, config = _require_loaded()
    data = collect_viz_data(model, tokenizer, prompt, config)
    template = _get_viz_template()
    json_str = json.dumps(data)
    inner_html = template.replace("{{VIZ_DATA_JSON}}", json_str)
    # Wrap in a full HTML document for the iframe
    doc = (
        '<!DOCTYPE html><html><head><meta charset="UTF-8">'
        '<style>body{margin:0;background:#0f172a;}</style></head>'
        f'<body>{inner_html}</body></html>'
    )
    # Escape for srcdoc attribute (HTML entity encoding)
    escaped = html_lib.escape(doc, quote=True)
    iframe = (
        f'<iframe srcdoc="{escaped}" '
        f'style="width:100%;height:750px;border:none;border-radius:8px;" '
        f'sandbox="allow-scripts"></iframe>'
    )
    return iframe


# ═══════════════════════════════════════════════════════════════
# Tab 8: Build Steps (12-step build tour)
# ═══════════════════════════════════════════════════════════════

# Step content: (title, body markdown, image-key or None)
# image-key resolves via build_viz.render_all() for structural diagrams,
# or via a small helper for the teach.py reused slides (step 2, 6, 8).

_BUILD_STEPS: list[tuple[str, str, str | None]] = [
    (
        "1 · Pin the hyperparameters  ·  `config.py`",
        "Start here, not at `model.py`. A single dataclass owning every tunable "
        "keeps the rest of the code uncluttered and reproducibility mechanical.\n\n"
        "**Non-obvious choice:** `vocab_size` starts at **0**. Every script that "
        "builds the model must set `config.vocab_size = tokenizer.vocab_size` "
        "*before* instantiation. This avoids hardcoding a vocab size and keeps "
        "the tokenizer and model coupled through a single variable.\n\n"
        "Derived values live in `@property`, with assertions that catch "
        "misconfiguration at access time rather than during a forward pass.",
        "config",
    ),
    (
        "2 · Byte-level BPE tokenizer  ·  `tokenizer.py`",
        "BPE from scratch in ~200 lines. Byte-level means the base vocab is "
        "always 256 (every possible byte), so any UTF-8 text round-trips "
        "losslessly — no `<UNK>` explosion on new characters.\n\n"
        "**Vocabulary layout:**\n"
        "```\n"
        "0-3     : <PAD>, <BOS>, <EOS>, <UNK>    (4 specials)\n"
        "4-259   : raw bytes 0x00-0xFF           (256 base)\n"
        "260+    : learned BPE merges\n"
        "```\n\n"
        "**Gotcha — test corpus variety:** A corpus of `\"hello world. \" * 100` "
        "collapses into a **single** token after ~10 merges. BPE tests need real "
        "variation. The visual below shows one prompt tokenized (this is slide 01 "
        "of the Train Reports tab, reused here).",
        "teach_01",  # rendered from teach.py on demand
    ),
    (
        "3 · Sliding-window dataset  ·  `dataset.py`",
        "Turns a token stream into `(input_ids, targets)` pairs where "
        "`targets = input_ids shifted +1`. This is the core causal-LM "
        "training invariant.\n\n"
        "**Default stride is `seq_len // 2`** (50% overlap). Non-overlap is a "
        "cleaner evaluation metric but halves the sample count; the overlap is "
        "slightly redundant for training but improves boundary next-token "
        "prediction.\n\n"
        "Tested by an invariant: `target[:-1] == input[1:]` — "
        "`tests/test_dataset.py::test_target_is_input_shifted_by_one`.",
        "sliding",
    ),
    (
        "4 · The Transformer block  ·  `model.py`",
        "Build bottom-up: **RMSNorm → RoPE → Attention → SwiGLU → Block**. "
        "Each component has a test in `tests/test_model.py`.\n\n"
        "**Pre-Norm** (normalize *before* the sublayer, add residual *after*) "
        "is what keeps deep stacks stable without a learning-rate-warmup "
        "knife-edge. Modern LLMs all use this pattern.\n\n"
        "**Weight tying:** `self.lm_head.weight = self.token_emb.weight` — "
        "not a copy, the SAME tensor. Test with `is` not `torch.equal`.\n\n"
        "**GPT-2 residual-projection init:** scale by `1/√(2·n_layers)` to keep "
        "activation variance stable as depth grows. Easy to miss; preserve it.",
        "block",
    ),
    (
        "5 · Training loop  ·  `train.py`",
        "The ingredients, in descending order of impact:\n\n"
        "1. **bf16 mixed precision** on Ampere+ — 2× speedup, basically free\n"
        "2. **AdamW** with weight-decay groups — decay on 2D weights only\n"
        "3. **Linear warmup → cosine decay** to 10% of peak LR (200 warmup steps)\n"
        "4. **Grad clip 1.0** — cheap insurance against loss spikes\n"
        "5. **Val split is sequential 90/10** — random split leaks via overlapping "
        "windows\n"
        "6. **\"Best\" checkpoint on *val* loss**, not train loss. Otherwise the "
        "\"best\" is usually right before overfitting starts.\n\n"
        "At the end, we dump `loss_curve.png` — students can SEE the gap between "
        "train and val open up. That IS overfitting.",
        "training",
    ),
    (
        "6 · Autoregressive generation (naïve)",
        "Forward the full context, sample from last-position logits, append, "
        "repeat. The sampling filter (temperature → top-k → top-p) is shared "
        "between `generate()`, `generate_fast()`, and `teach.py`'s sampling "
        "rollout — one source of truth, `_sample_from_logits()`.\n\n"
        "The `scatter` at the end of top-p looks fragile but is correct: "
        "`sorted_idx` is a permutation over ALL vocab positions, so every slot "
        "is written.\n\n"
        "**This is slide 11 of the Train Reports tab** — watch 10 decode steps with the "
        "top-5 candidates at each step and which one was sampled.",
        "teach_11",  # sampling rollout
    ),
    (
        "7 · KV cache  ·  `generate_fast()`",
        "Without a cache, every decode step recomputes Q, K, V for the entire "
        "context. With a cache, you keep K, V from prior steps and only compute "
        "them for the **new token**. Each step becomes "
        "`O(total_len)` instead of `O(total_len²)`.\n\n"
        "**Two invariants that are easy to miss:**\n\n"
        "- **RoPE offset.** The new token's RoPE angle must match what it would "
        "have been in a single-pass forward — so `start_pos = past_len`, not 0.\n"
        "- **Causal mask.** Applied during **prefill** (T>1), skipped during "
        "**decode** (T=1).\n\n"
        "**Correctness test:** `prefill(prompt)` then 1-step decode must match "
        "a full forward on `prompt`. `test_kv_cache_matches_full_pass` asserts "
        "`max |Δlogit| < 1e-4`. This is the test that catches a broken RoPE "
        "offset.\n\n"
        "**Keep both paths.** Deleting `generate()` loses the pedagogical A/B "
        "comparison. The Benchmark tab literally shows both numbers side by side.",
        "kv_cache",
    ),
    (
        "8 · Teaching hooks  ·  `teach.py`",
        "Pattern worth stealing: capture intermediate tensors via "
        "**forward hooks** without modifying `model.py` at all.\n\n"
        "```python\n"
        "block.attn.register_forward_pre_hook(attn_hook)\n"
        "def attn_hook(module, inp):\n"
        "    x = inp[0]\n"
        "    # Recompute Q, K, V, scores, weights from x\n"
        "    store['attn_weights'] = ...\n"
        "```\n\n"
        "The alternative — sprinkling `return_attention=True` kwargs through "
        "the production code — poisons `model.py`. Hooks keep teaching "
        "instrumentation strictly orthogonal to the model.\n\n"
        "**Re-render on-the-fly** (don't cache PNGs to disk and serve those): "
        "the whole point is that changing the prompt flips the attention. "
        "That's why this UI calls `teach.py` functions directly.\n\n"
        "The visual below is slide 08 — all attention heads of one layer "
        "side-by-side, each learning a different pattern.",
        "teach_08",  # all heads
    ),
    (
        "9 · Attention visualisation  ·  `visualise.py`",
        "Two things worth rendering:\n\n"
        "1. **Per-head heatmap** — `attn_weights[layer][0, head]`, shape `(T, T)`, "
        "viridis colormap, row = query, column = key.\n"
        "2. **Attention rollout** (Abnar & Zuidema 2020) — multiply attention "
        "matrices across all layers to see *effective* attention flow. More "
        "interpretable than any single layer.\n\n"
        "Both are in the **Attention** tab. Pick any prompt + layer + head to "
        "see the heatmap live.",
        None,  # use the Attention tab directly
    ),
    (
        "10 · Gradio web UI  ·  `app.py`  (this page!)",
        "A single browser app entrypoint. All handlers import "
        "and call existing functions from `teach.py` and `visualise.py` — no "
        "logic re-implemented.\n\n"
        "**Two Gradio patterns worth learning:**\n\n"
        "1. **Streaming via `yield`.** A generator handler gives you per-token "
        "typewriter output for free.\n"
        "2. **Module-level singletons + `monkeypatch` for tests.** Load the "
        "model once at startup; tests inject a tiny model via "
        "`monkeypatch.setattr(app, \"_MODEL\", tiny_model)`. No I/O, "
        "sub-second UI smoke tests.\n\n"
        "**Version pin:** `gradio>=5,<6`. 4.44.x has a schema-introspection "
        "bug that crashes `launch()`. Upgrade fixes it.",
        "ui_layout",
    ),
    (
        "11 · Test suite  ·  `tests/`",
        "48 tests (42 unit + 6 UI smoke), CPU-only, ~18 s end-to-end. The "
        "highest-value tests are the ones that guard **invariants** rather "
        "than shapes:\n\n"
        "- `test_kv_cache_matches_full_pass` → RoPE `start_pos` invariant\n"
        "- `test_causal_mask_no_future_leakage` → perturbing token T-1 "
        "leaves earlier logits unchanged\n"
        "- `test_weight_tying` → `is`-check: `lm_head.weight` and "
        "`token_emb.weight` are the SAME tensor\n"
        "- `test_training_step_reduces_loss_on_tiny_overfit` → end-to-end "
        "backward/optimizer actually moves parameters\n\n"
        "Invariant tests survive refactoring. Shape tests get rewritten every "
        "time somebody touches the code.",
        "test_matrix",
    ),
    (
        "12 · Ship it  ·  `run.sh`",
        "Commands unified through one bash dispatcher:\n\n"
        "```bash\n"
        "bash run.sh setup       # venv + torch + TinyStories\n"
        "bash run.sh verify      # model.py __main__ sanity check\n"
        "bash run.sh train       # full training (~5-15 min on RTX 4080)\n"
        "bash run.sh generate --fast\n"
        "bash run.sh benchmark   # generate vs generate_fast\n"
        "bash run.sh teach       # 16 static PNGs\n"
        "bash run.sh visualise   # attention heatmaps + rollout\n"
        "bash run.sh ui          # this console\n"
        "bash run.sh test        # pytest suite\n"
        "```\n\n"
        "Explicit subcommands > magic auto-detection. "
        "A new user running `bash run.sh` with no arg gets a helpful usage message.\n\n"
        "**Further reading:** `BUILDING.md` (this page, as markdown), "
        "`REFERENCES.md` (slide ↔ Raschka book), "
        "`nanollm-guide.md` (RNN → Transformer theory).",
        None,
    ),
]


_BUILD_IMG_CACHE: dict[str, str] = {}


def _render_build_image(key: str) -> str | None:
    """Resolve a step's image-key to a PNG path. Idempotent & cached."""
    if key is None:
        return None
    if key in _BUILD_IMG_CACHE:
        return _BUILD_IMG_CACHE[key]

    outdir = tempfile.mkdtemp(prefix="nanollm_build_") \
        if "_BUILD_DIR" not in globals() else globals()["_BUILD_DIR"]
    globals()["_BUILD_DIR"] = outdir

    # Structural diagrams from build_viz
    diagram_paths = build_viz.render_all(outdir)
    if key in diagram_paths:
        _BUILD_IMG_CACHE[key] = diagram_paths[key]
        return diagram_paths[key]

    # Reused slides from teach.py (rendered with a fixed demo prompt)
    if key.startswith("teach_"):
        path = _render_teach_reuse(key, outdir)
        if path is not None:
            _BUILD_IMG_CACHE[key] = path
        return path

    return None


def _render_teach_reuse(key: str, outdir: str) -> str | None:
    """Render one of the teach.py slides using a fixed demo prompt."""
    try:
        model, tokenizer, config = _require_loaded()
    except RuntimeError:
        return None

    demo_prompt = "The cat sat on the"
    token_ids = tokenizer.encode(demo_prompt, add_special=False)
    if len(token_ids) < 2:
        return None
    if len(token_ids) > config.max_seq_len:
        token_ids = token_ids[:config.max_seq_len]
    labels = token_labels(tokenizer, token_ids, max_len=8)
    idx = torch.tensor([token_ids], device=config.device)

    plt_ = _get_plt()

    # Run capture once; reuse for all teach_* slides in this session
    if "_BUILD_CAPTURE" not in globals():
        cap = ForwardCapture(model, target_layer=0)
        with torch.no_grad():
            model(idx)
        cap.remove()
        globals()["_BUILD_CAPTURE"] = cap
        globals()["_BUILD_LABELS"] = labels
        globals()["_BUILD_DEMO_IDS"] = token_ids
    cap = globals()["_BUILD_CAPTURE"]
    labels = globals()["_BUILD_LABELS"]
    token_ids = globals()["_BUILD_DEMO_IDS"]

    path = os.path.join(outdir, f"{key}.png")

    if key == "teach_01":
        slide_01_tokenization(plt_, tokenizer, demo_prompt, token_ids, path)
    elif key == "teach_08":
        slide_08_all_heads(plt_, cap.store, labels, path)
    elif key == "teach_11":
        slide_11_sampling_rollout(plt_, model, tokenizer, idx, config.device,
                                  0.8, 40, 0.9, path)
    else:
        return None

    return path


def render_step_panel(step_idx: int) -> tuple[str, str | None]:
    """Return (markdown, image_path) for a given step index."""
    if not (0 <= step_idx < len(_BUILD_STEPS)):
        return "Unknown step", None
    title, body, img_key = _BUILD_STEPS[step_idx]
    md = f"### {title}\n\n{body}"
    img = _render_build_image(img_key) if img_key else None
    return md, img


# ═══════════════════════════════════════════════════════════════
# Notebook-derived tabs (Tokenizer / Dataset / Components / KV Cache)
#
# Mirror sections 2, 3, 4, and 9 of LLM101_From_Scratch.ipynb. The
# pure-render helpers live in notebook_viz.py; the handlers below are
# thin glue that writes the returned Figures to a tempdir.
# ═══════════════════════════════════════════════════════════════

def _save_fig(fig, prefix: str = "nanollm_viz_") -> str:
    """Write a matplotlib Figure to a fresh tempdir and return the PNG path."""
    outdir = tempfile.mkdtemp(prefix=prefix)
    path = os.path.join(outdir, "viz.png")
    fig.savefig(path, dpi=130)
    plt.close(fig)
    return path


def run_tokenizer(text: str):
    """Encode user text → breakdown text + vocab/compression overview chart."""
    _, tokenizer, _ = _require_loaded()
    text = text or "Hello world!"
    breakdown = nv.encode_breakdown(tokenizer, text)
    return nv.format_breakdown(breakdown), _save_fig(
        nv.draw_tokenizer_overview(tokenizer), "nanollm_tok_"
    )


def run_dataset(text: str, seq_len: int, stride: int, sample_idx: int):
    """Tokenize the input text and render the input/target shift + windows."""
    _, tokenizer, config = _require_loaded()
    text = text or ("Once upon a time, a little cat sat on the mat. " * 20)
    tokens = tokenizer.encode(text, add_special=False)
    if not tokens:
        return None, "Empty tokenization."
    seq_len = max(2, min(int(seq_len), config.max_seq_len, max(2, len(tokens) - 1)))
    stride = max(1, int(stride))
    sample_idx = max(0, int(sample_idx))
    fig = nv.draw_window_view(
        tokenizer, tokens, seq_len=seq_len, stride=stride,
        sample_idx=sample_idx, n_show=min(12, seq_len),
    )
    summary = (f"Tokenized: {len(text):,} chars → {len(tokens):,} tokens · "
               f"seq_len={seq_len}, stride={stride}, sample_idx={sample_idx}")
    return _save_fig(fig, "nanollm_ds_"), summary


def run_component(name: str):
    """Render the (text, figure) pair for one TransformerBlock sub-tab."""
    _, _, config = _require_loaded()
    if name == "rmsnorm":
        return nv.rmsnorm_summary(config), _save_fig(
            nv.draw_rmsnorm_dist(config), "nanollm_rmsnorm_")
    if name == "rope":
        return nv.rope_summary(config), _save_fig(
            nv.draw_rope_demo(config), "nanollm_rope_")
    if name == "attention":
        return nv.attention_summary(config), _save_fig(
            nv.draw_causal_mask(config), "nanollm_attn_")
    if name == "swiglu":
        return nv.swiglu_summary(config), _save_fig(
            nv.draw_swiglu_breakdown(config), "nanollm_swiglu_")
    raise ValueError(f"Unknown component: {name}")


def run_kv_cache(seq_len: int, multi_T: int, sweep_str: str, gen_len: int):
    """§9 deep-dive: single-step + multi-step equivalence + length-sweep timing."""
    model, _, config = _require_loaded()
    seq_len = int(seq_len)
    multi_T = int(multi_T)
    gen_len = int(gen_len)

    single = nv.format_kv_single_step(
        nv.kv_cache_single_step(model, config, T=seq_len)
    )
    multi = nv.format_kv_multi_step(
        nv.kv_cache_multi_step(model, config, T=multi_T)
    )

    try:
        prompt_lens = [int(p.strip()) for p in sweep_str.split(",") if p.strip()]
    except ValueError:
        prompt_lens = []
    if not prompt_lens:
        prompt_lens = [5, 10, 20, 40]

    fig, sweep_text = nv.draw_length_sweep(model, config, prompt_lens,
                                           gen_len=gen_len)
    return single, multi, _save_fig(fig, "nanollm_kv_sweep_"), sweep_text


# ═══════════════════════════════════════════════════════════════
# UI layout
# ═══════════════════════════════════════════════════════════════

def build_ui() -> gr.Blocks:
    """Assemble the 4-tab Gradio app. Returns the Blocks object (unlaunched)."""
    # Determine UI ranges from the loaded config
    n_layers = _CONFIG.n_layers if _CONFIG is not None else 6
    n_heads = _CONFIG.n_heads if _CONFIG is not None else 6
    max_seq = _CONFIG.max_seq_len if _CONFIG is not None else 256

    _css = "#sidebar-nav .wrap { display: flex !important; flex-direction: column !important; gap: 4px !important; }"
    with gr.Blocks(title="LLM101 - Console",
                   theme=gr.themes.Soft(primary_hue="blue"),
                   css=_css) as demo:
        gr.Markdown(
            "# LLM101 - Console\n"
            "A ~15M-parameter GPT-style decoder-only transformer, built from scratch in PyTorch. "
            "📓 [Open in JupyterLab](http://127.0.0.1:8888/lab/tree/LLM101_From_Scratch.ipynb) "
            "— run `bash run.sh notebook` in a terminal first to start the server.\n\n"
            "- **Tokenizer** (§2) — BPE vocabulary: how raw text is split into tokens\n"
            "- **Dataset** (§3) — Sliding-window sampling: input/target pairs for training\n"
            "- **TransformerBlock** (§4) — Block anatomy: RMSNorm · RoPE · Causal Attention · SwiGLU FFN\n"
            "- **Train** (§6) — Live training: adjust hyperparameters, watch loss fall in real time\n"
            "- **Fine-tune** — Adapt the pre-trained model on your own text at a lower learning rate\n"
            "- **Train Reports** (§5) — 16-slide forward-pass walkthrough: every tensor transformation visualised\n"
            "- **Attention** (§8) — Per-head heatmap + attention rollout for any prompt\n"
            "- **Visualize** (§8) — Animated 3-panel view: hidden states, norms, and attention across layers\n"
            "- **Generate** (§7) — Token-by-token generation with temperature / top-k / top-p controls\n"
            "- **Benchmark** (§9) — Throughput: `generate()` (no cache) vs `generate_fast()` (KV cache)\n"
            "- **KV Cache** (§9) — Correctness tests: single-step, multi-step equivalence, length-sweep timing\n"
            "- **Architecture** — Interactive Mermaid flowchart of the full end-to-end pipeline"
        )
        gr.Markdown(_STATUS)

        _NAV_CHOICES = [
            "Tokenizer", "Dataset", "TransformerBlock", "Train",
            "Fine-tune", "Train Reports", "Attention", "Visualize", "Generate",
            "Benchmark", "KV Cache", "Architecture",
        ]

        with gr.Sidebar(open=True):
            gr.Markdown("### LLM101")
            _nav = gr.Radio(
                choices=_NAV_CHOICES,
                value="Tokenizer",
                label="",
                elem_id="sidebar-nav",
            )

        # ── hidden legacy panels (not in nav) ──
        with gr.Column(visible=False):
            gr.Markdown(
                "A 12-step tour of how this project was actually built, in the order "
                "a new builder should follow. Complements the standalone `BUILDING.md` "
                "with inline visualisations. Use ◀ / ▶ to walk through, or pick any "
                "step from the list."
            )
            with gr.Row():
                with gr.Column(scale=1, min_width=220):
                    build_radio = gr.Radio(
                        choices=[s[0].split("·")[0].strip() for s in _BUILD_STEPS],
                        value=_BUILD_STEPS[0][0].split("·")[0].strip(),
                        label="Step", interactive=True,
                    )
                    with gr.Row():
                        build_prev = gr.Button("◀ Previous", size="sm")
                        build_next = gr.Button("Next ▶", size="sm", variant="primary")
                with gr.Column(scale=3):
                    build_md = gr.Markdown(_BUILD_STEPS[0][1])  # eager first step
                    build_img = gr.Image(
                        label="Visualisation", type="filepath",
                        show_label=False, height=500,
                    )

            def _idx_from_choice(choice: str) -> int:
                for i, s in enumerate(_BUILD_STEPS):
                    if s[0].split("·")[0].strip() == choice:
                        return i
                return 0

            def _img_or_hidden(img):
                """Return the image or hide the component if None."""
                if img is None:
                    return gr.update(value=None, visible=False)
                return gr.update(value=img, visible=True)

            def on_select(choice):
                md, img = render_step_panel(_idx_from_choice(choice))
                return md, _img_or_hidden(img)

            def on_prev(choice):
                i = max(0, _idx_from_choice(choice) - 1)
                new_choice = _BUILD_STEPS[i][0].split("·")[0].strip()
                md, img = render_step_panel(i)
                return new_choice, md, _img_or_hidden(img)

            def on_next(choice):
                i = min(len(_BUILD_STEPS) - 1, _idx_from_choice(choice) + 1)
                new_choice = _BUILD_STEPS[i][0].split("·")[0].strip()
                md, img = render_step_panel(i)
                return new_choice, md, _img_or_hidden(img)

            build_radio.change(on_select, inputs=build_radio,
                               outputs=[build_md, build_img])
            build_prev.click(on_prev, inputs=build_radio,
                             outputs=[build_radio, build_md, build_img])
            build_next.click(on_next, inputs=build_radio,
                             outputs=[build_radio, build_md, build_img])

            # Eager-render the first step's image on load so the initial view
            # isn't empty
            demo.load(
                lambda: on_select(_BUILD_STEPS[0][0].split("·")[0].strip()),
                inputs=None,
                outputs=[build_md, build_img],
            )

        # ── Panel: Tokenizer ──
        with gr.Column(visible=True) as _panel_tokenizer:
            gr.Markdown(
                "### Byte-level BPE — see exactly how text becomes tokens\n\n"
                "Mirrors **§2** of `LLM101_From_Scratch.ipynb`. Type any text "
                "and see the encode/decode round-trip with a per-token kind "
                "label (SPECIAL · BYTE · MERGE), plus the vocabulary "
                "composition and the BPE compression ratio across sample texts."
            )
            with gr.Row():
                with gr.Column(scale=1):
                    tok_in = gr.Textbox(
                        label="Text to tokenize",
                        value="To be, or not to be, that is the question.",
                        lines=3,
                    )
                    tok_btn = gr.Button("Tokenize", variant="primary")
                    tok_breakdown = gr.Code(
                        label="Round-trip breakdown",
                        language=None, interactive=False,
                    )
                with gr.Column(scale=1):
                    tok_overview = gr.Image(
                        label="Vocabulary breakdown + compression",
                        type="filepath",
                    )
                    gr.Markdown(
                        "**Left:** The vocabulary has 3 tiers — 4 special control tokens "
                        "(PAD/BOS/EOS/UNK), 256 byte-level fallback tokens that can encode "
                        "*any* Unicode character losslessly, and the bulk: BPE merge tokens "
                        "learned from the corpus that represent common byte sequences as a "
                        "single unit.\n\n"
                        "**Right:** Compression ratio = bytes in the original text ÷ number "
                        "of tokens produced. A ratio of 4× means 4 bytes were packed into "
                        "1 token. The red dashed line at 1.0 is the no-compression baseline "
                        "(one token per byte). Richer, repetitive text compresses better "
                        "because BPE has seen those patterns during training."
                    )
            tok_btn.click(run_tokenizer,
                          inputs=[tok_in],
                          outputs=[tok_breakdown, tok_overview])
            demo.load(run_tokenizer, inputs=[tok_in],
                      outputs=[tok_breakdown, tok_overview])

        with gr.Column(visible=False) as _panel_dataset:
            gr.Markdown(
                "### Sliding-window dataset — input vs target shift, window overlap\n\n"
                "Mirrors **§3** of the notebook. The corpus is encoded once, "
                "then sliced into overlapping windows of length `seq_len` with "
                "the chosen `stride`. For causal LM, **target = input shifted "
                "right by 1**: at every position the model predicts the *next* "
                "token from everything to its left."
            )
            with gr.Row():
                with gr.Column(scale=1):
                    ds_text = gr.Textbox(
                        label="Source text (will be tokenized)",
                        value=("Once upon a time, a little cat sat on the mat. "
                               "The dog watched, then chased a bird. "
                               "Tomorrow we will go to the park. " * 8),
                        lines=4,
                    )
                    ds_seq = gr.Slider(
                        4, max_seq, value=min(32, max_seq), step=4,
                        label="seq_len (window length)",
                    )
                    ds_stride = gr.Slider(
                        1, max_seq, value=min(16, max_seq), step=1,
                        label="stride (offset between windows)",
                    )
                    ds_idx = gr.Slider(
                        0, 8, value=0, step=1,
                        label="sample index (which window to inspect)",
                    )
                    ds_btn = gr.Button("Render windows", variant="primary")
                    ds_status = gr.Markdown()
                with gr.Column(scale=2):
                    ds_img = gr.Image(label="Shift + windows", type="filepath")
                    gr.Markdown(
                        "<small>"
                        "**seq_len** — how many tokens the model sees at once (its context window). "
                        "Larger values give richer context but consume more GPU memory and require longer sequences in the training data.<br>"
                        "**stride** — how far the window slides between consecutive samples. "
                        "stride &lt; seq_len creates overlap so tokens near a boundary appear in multiple training windows, seen in different contexts.<br>"
                        "**sample index** — selects which overlapping window to display in the top panel. "
                        "Increasing by 1 shifts the highlighted window right by <code>stride</code> tokens — visible in the Gantt chart above."
                        "</small>"
                    )
            ds_btn.click(run_dataset,
                         inputs=[ds_text, ds_seq, ds_stride, ds_idx],
                         outputs=[ds_img, ds_status])
            demo.load(run_dataset,
                      inputs=[ds_text, ds_seq, ds_stride, ds_idx],
                      outputs=[ds_img, ds_status])

            with gr.Accordion("data/corpus.txt — preview", open=False):
                corpus_stats = gr.Markdown()
                corpus_preview = gr.Textbox(
                    label="First 3 000 characters",
                    lines=16, max_lines=16,
                    interactive=False, show_copy_button=True,
                )
                def _load_corpus():
                    path = "data/corpus.txt"
                    if not os.path.exists(path):
                        return "_corpus not found — run `bash run.sh setup` first_", ""
                    size = os.path.getsize(path)
                    with open(path, encoding="utf-8", errors="replace") as f:
                        text = f.read()
                    lines = text.count("\n")
                    preview = text[:3000] + ("\n…" if len(text) > 3000 else "")
                    stats = (f"**{path}** · {size:,} bytes · {len(text):,} chars "
                             f"· {lines:,} lines · {len(text.split()):,} words")
                    return stats, preview
                demo.load(_load_corpus, inputs=None, outputs=[corpus_stats, corpus_preview])

        with gr.Column(visible=False) as _panel_tb:
            gr.Markdown(
                "### Atomic components — RMSNorm · RoPE · Attention · SwiGLU\n\n"
                "Mirrors **§4** of the notebook. Each sub-tab instantiates "
                "**one** sub-module against the current config and shows what "
                "it does in isolation, before they're stacked into the full "
                "TransformerBlock. Numbers update if you re-train with a "
                "different config."
            )
            with gr.Row():
                # ── Left column: architecture diagram ──
                with gr.Column(scale=1):
                    gr.Image(
                        value=_save_fig(nv.draw_transformer_block_diagram(), "comp_diagram"),
                        label="TransformerBlock — how the components connect",
                        type="filepath", show_download_button=False,
                        height=800,
                    )

                # ── Right column: component detail tabs ──
                with gr.Column(scale=1):
                    with gr.Tabs():
                        with gr.Tab("RMSNorm"):
                            rms_txt = gr.Code(label="Summary", language=None,
                                              interactive=False, lines=18, max_lines=18)
                            gr.Markdown(
                                "<small><b>Summary:</b> RMSNorm rescales each activation vector so its root-mean-square equals 1, "
                                "then multiplies by a learnable per-dimension weight γ. Unlike LayerNorm it skips the mean-subtraction step, "
                                "saving 384 parameters per layer (50% reduction) with no measurable quality loss.</small>"
                            )
                            gr.HTML("<hr style='border:none;border-top:1px solid #444;margin:6px 0'>")
                            gr.Markdown(
                                "<small><b>Chart:</b> The blue histogram (Before) is wide and off-centre — raw activations can grow large and destabilise gradients. "
                                "The green histogram (After) is tightly concentrated around 0 with std ≈ 1, giving the next layer consistent input scale regardless of sequence content.</small>"
                            )
                            rms_img = gr.Image(
                                label="Activation distribution before / after",
                                type="filepath", height=280,
                            )

                        with gr.Tab("RoPE"):
                            rope_txt = gr.Code(label="Summary", language=None,
                                               interactive=False, lines=18, max_lines=18)
                            gr.Markdown(
                                "<small><b>Summary:</b> RoPE encodes position by rotating each query/key vector in 2D planes using sine and cosine functions of the position index. "
                                "The rotation angle grows with frequency index, so nearby tokens share similar angles while distant tokens diverge — "
                                "this makes attention scores naturally decay with distance without any learned positional embedding table.</small>"
                            )
                            gr.HTML("<hr style='border:none;border-top:1px solid #444;margin:6px 0'>")
                            gr.Markdown(
                                "<small><b>Chart:</b> Left two panels show the cos and sin frequency tables across positions (rows) and dimension pairs (columns) — "
                                "each column oscillates at a different frequency. The right panel shows how a single token's vector rotates as its position increases; "
                                "the coloured arrows trace the unit circle, one arrow per position.</small>"
                            )
                            rope_img = gr.Image(
                                label="cos / sin tables + unit-vector rotation",
                                type="filepath", height=280,
                            )

                        with gr.Tab("Attention"):
                            cmask_txt = gr.Code(label="Summary", language=None,
                                                interactive=False, lines=18, max_lines=18)
                            gr.Markdown(
                                "<small><b>Summary:</b> The causal mask enforces the autoregressive constraint: position i may only attend to positions ≤ i. "
                                "It is implemented as an upper-triangular matrix of −∞ values added to the raw attention scores before softmax, "
                                "driving those weights to zero so future tokens contribute nothing to the current prediction.</small>"
                            )
                            gr.HTML("<hr style='border:none;border-top:1px solid #444;margin:6px 0'>")
                            gr.Markdown(
                                "<small><b>Chart:</b> Dark cells (1 = visible) form the lower triangle including the diagonal — a token always attends to itself and all earlier tokens. "
                                "White cells (0 = masked) are the future positions that get −∞ before softmax. "
                                "The strict triangular shape is why the model can be trained on all positions in parallel yet still learns left-to-right generation.</small>"
                            )
                            cmask_img = gr.Image(
                                label="Causal mask (1 = visible, 0 = masked)",
                                type="filepath", height=280,
                            )

                        with gr.Tab("SwiGLU"):
                            sw_txt = gr.Code(label="Summary", language=None,
                                             interactive=False, lines=18, max_lines=18)
                            gr.Markdown(
                                "<small><b>Summary:</b> SwiGLU splits the FFN expansion into two parallel projections (gate and value). "
                                "The gate branch passes through Swish (x·σ(x)), producing a smooth gating signal that multiplies element-wise with the value branch — "
                                "selectively suppressing or amplifying each hidden dimension before the final down-projection. "
                                "This gives the network a multiplicative interaction that ReLU FFNs lack, improving expressiveness for the same parameter count.</small>"
                            )
                            gr.HTML("<hr style='border:none;border-top:1px solid #444;margin:6px 0'>")
                            gr.Markdown(
                                "<small><b>Chart:</b> The three bars show the parameter count of each linear layer: gate_proj, up_proj (both d_model → d_ff), and down_proj (d_ff → d_model). "
                                "All three are the same size; together they account for roughly half the model's total parameters.</small>"
                            )
                            sw_img = gr.Image(
                                label="Parameter breakdown across 3 projections",
                                type="filepath", height=280,
                            )

            demo.load(lambda: run_component("rmsnorm"),
                      inputs=None, outputs=[rms_txt, rms_img])
            demo.load(lambda: run_component("rope"),
                      inputs=None, outputs=[rope_txt, rope_img])
            demo.load(lambda: run_component("attention"),
                      inputs=None, outputs=[cmask_txt, cmask_img])
            demo.load(lambda: run_component("swiglu"),
                      inputs=None, outputs=[sw_txt, sw_img])

        with gr.Column(visible=False) as _panel_train:
            gr.Markdown(
                "**Trigger training from the UI and watch progress live.** "
                "Training writes `checkpoints/best.pt` and `checkpoints/loss_curve.png`. "
                "When it finishes, the Generate / Train Reports / Attention tabs pick up the "
                "new model automatically — no restart. "
                "Needs `data/corpus.txt` (run `bash run.sh setup` first)."
            )
            with gr.Row():
                # ── Left: controls ──
                with gr.Column(scale=1, min_width=280):
                    train_epochs = gr.Slider(
                        1, 30, value=10, step=1, label="max_epochs",
                    )
                    train_batch = gr.Slider(
                        8, 128, value=64, step=8, label="batch_size",
                    )
                    train_lr = gr.Slider(
                        1e-5, 1e-3, value=3e-4, step=1e-5, label="learning_rate",
                    )
                    train_dropout = gr.Slider(
                        0.0, 0.4, value=0.1, step=0.05, label="dropout",
                    )
                    train_warmup = gr.Slider(
                        0, 500, value=0, step=10, label="warmup_steps",
                    )
                    with gr.Row():
                        train_btn  = gr.Button("Start",   variant="primary",   scale=2)
                        train_stop_btn    = gr.Button("Stop",    variant="stop",      scale=1)
                        train_restart_btn = gr.Button("Restart", variant="secondary", scale=1)
                    gr.HTML("<hr style='border:none;border-top:1px solid #444;margin:8px 0'>")
                    gr.Markdown(
                        "<small>"
                        "**max_epochs** — Full passes through the dataset. Too few = underfit, too many = overfit; watch val_loss to find the sweet spot.<br><br>"
                        "**batch_size** — Sequences processed in parallel per step. Larger = smoother gradients but more memory; lower (8–16) if you hit OOM.<br><br>"
                        "**learning_rate** — Peak LR after warmup. Too high = unstable, too low = slow convergence; 3e-4 is the AdamW standard.<br><br>"
                        "**dropout** — Fraction of activations zeroed during training. Regularisation that prevents overfitting — the model can't rely on any single neuron.<br><br>"
                        "**warmup_steps** — Steps where LR ramps from 0 to peak before cosine decay. Prevents early gradient explosion; 0 = auto-scale to 10% of total steps."
                        "</small>"
                    )
                # ── Right: outputs ──
                with gr.Column(scale=2):
                    train_log = gr.Textbox(
                        label="Training log (streaming)",
                        value="(press Start to begin)",
                        lines=14, max_lines=20, show_copy_button=True,
                        interactive=False,
                    )
                    train_plot = gr.Image(
                        label="Loss curve (live)", type="filepath", height=300,
                    )
                    gr.Markdown(
                        "<small>Blue = training loss (per step), orange = validation loss (per epoch). "
                        "Both should fall over time. A rising val loss while train loss keeps falling = overfitting. "
                        "The gap between them shows how well the model generalises beyond the training data.</small>"
                    )
                    train_status = gr.Textbox(
                        label="Status", value="idle", interactive=False,
                    )
                    gr.Markdown(
                        "<small><b>idle</b> — no run yet · <b>running</b> — training in background (safe to switch tabs) · "
                        "<b>done</b> — finished, model reloaded · <b>error</b> — check log above.</small>"
                    )
                    train_samples = gr.Textbox(
                        label="Generation samples (newest first)",
                        value="(samples appear here after each epoch completes)",
                        lines=8, max_lines=20, interactive=False,
                    )
                    gr.Markdown(
                        "<small>After each epoch the model generates from a fixed prompt — newest epoch on top. "
                        "Early epochs produce random-looking text; by epoch 3–5 you should see coherent words and simple sentences. "
                        "Scroll down to compare earlier epochs and track how fluency improves over time.</small>"
                    )

            _train_inputs = [train_epochs, train_batch, train_lr,
                             train_dropout, train_warmup]
            train_btn.click(
                train_start, inputs=_train_inputs, outputs=[train_status],
            )
            train_stop_btn.click(
                train_stop, inputs=[], outputs=[train_status],
            )
            train_restart_btn.click(
                train_restart, inputs=_train_inputs, outputs=[train_status],
            )
            # Poll training state every 2 s — works even when tab is not visible
            train_timer = gr.Timer(value=2)
            train_timer.tick(poll_train,
                             outputs=[train_log, train_plot, train_status, train_samples])

        with gr.Column(visible=False) as _panel_finetune:
            gr.Markdown(
                "### Fine-tuning — specialise the pre-trained model on your own text\n\n"
                "**What is fine-tuning?**  The pre-trained model has already learned "
                "general language patterns from TinyStories (~1.5 M characters). "
                "Fine-tuning continues training from those learned weights on a small "
                "custom corpus, so the model adapts its style, vocabulary, and topics "
                "while retaining its general language ability.\n\n"
                "**Why a lower learning rate?**  The model's weights are already in a "
                "good region of parameter space. A large step (like 3e-4 used for "
                "pre-training) would destroy that structure — the model would "
                "'forget' general language and overfit the tiny corpus. "
                "A smaller LR (1e-5 – 1e-4) nudges the weights gently.\n\n"
                "**What to paste:** Any text you want the model to imitate — "
                "a style guide, a short story, song lyrics, domain-specific prose. "
                "Aim for **at least 500 characters** (a few paragraphs). "
                "The tokenizer is reused from pre-training — no vocabulary changes.\n\n"
                "**Output:** `checkpoints/finetuned.pt` — the Generate tab will "
                "switch to this model automatically when fine-tuning finishes."
            )
            with gr.Row():
                # ── Left: controls ──
                with gr.Column(scale=1, min_width=280):
                    ft_text = gr.Textbox(
                        label="Custom corpus (paste your text here)",
                        placeholder=(
                            "Paste at least 200 characters of text here.\n\n"
                            "Example: a short story, article, code comments, "
                            "domain-specific prose — anything you want the model "
                            "to learn to generate in that style."
                        ),
                        lines=10, max_lines=20,
                    )
                    ft_ckpt = gr.Textbox(
                        label="Base checkpoint",
                        value="checkpoints/best.pt",
                        info="Path to the pre-trained weights to start from.",
                    )
                    ft_lr = gr.Slider(
                        1e-5, 1e-4, value=5e-5, step=1e-6,
                        label="learning_rate",
                        info="Keep well below the pre-training LR (3e-4) to avoid catastrophic forgetting.",
                    )
                    ft_epochs = gr.Slider(
                        1, 20, value=5, step=1, label="max_epochs",
                        info="More epochs = more overfitting on small corpora. 3–5 is usually enough.",
                    )
                    ft_batch = gr.Slider(
                        4, 64, value=16, step=4, label="batch_size",
                        info="Smaller batches work better for tiny corpora; 8–16 is typical.",
                    )
                    with gr.Row():
                        ft_btn  = gr.Button("Start fine-tuning", variant="primary", scale=2)
                        ft_stop_btn = gr.Button("Stop", variant="stop", scale=1)
                    gr.HTML("<hr style='border:none;border-top:1px solid #444;margin:8px 0'>")
                    gr.Markdown(
                        "<small>"
                        "**learning_rate** — use 5e-5 as a safe default. "
                        "Pre-training used 3e-4; going higher risks 'catastrophic forgetting' "
                        "where the model overwrites general language ability with the new corpus.<br><br>"
                        "**max_epochs** — with a small corpus (< 5 000 tokens) the model can overfit "
                        "in 3–5 epochs. Watch val_loss: if it starts rising while train_loss falls, "
                        "stop early (or reduce epochs).<br><br>"
                        "**batch_size** — small corpus = fewer windows = smaller optimal batch. "
                        "If you see 'not enough tokens' warnings, lower this.<br><br>"
                        "**What changes and what doesn't** — only the weight *values* change. "
                        "Architecture (d_model, n_layers, n_heads), tokenizer, and max_seq_len "
                        "are all frozen at the pre-trained checkpoint's settings."
                        "</small>"
                    )
                # ── Right: outputs ──
                with gr.Column(scale=2):
                    ft_log = gr.Textbox(
                        label="Fine-tuning log",
                        value="(press Start to begin)",
                        lines=12, max_lines=20, show_copy_button=True,
                        interactive=False,
                    )
                    ft_plot = gr.Image(
                        label="Loss curve (live)", type="filepath", height=280,
                    )
                    gr.Markdown(
                        "<small>Blue = training loss, orange = validation loss. "
                        "Val loss higher than train loss is normal and expected for small corpora. "
                        "Stop if val loss starts diverging sharply (overfitting).</small>"
                    )
                    ft_status = gr.Textbox(
                        label="Status", value="idle", interactive=False,
                    )
                    ft_samples = gr.Textbox(
                        label="Generation samples (newest first)",
                        value="(samples appear here after each epoch)",
                        lines=6, max_lines=16, interactive=False,
                    )
                    gr.Markdown(
                        "<small>After each epoch the model generates from a fixed prompt using the "
                        "fine-tuned weights. Compare early vs late epochs to see style adaptation. "
                        "The final fine-tuned model is saved to "
                        "`checkpoints/finetuned.pt` and loaded into the Generate tab.</small>"
                    )

            ft_btn.click(
                ft_start,
                inputs=[ft_text, ft_ckpt, ft_lr, ft_epochs, ft_batch],
                outputs=[ft_status],
            )
            ft_stop_btn.click(ft_stop, inputs=[], outputs=[ft_status])
            ft_timer = gr.Timer(value=2)
            ft_timer.tick(poll_ft, outputs=[ft_log, ft_plot, ft_status, ft_samples])

        with gr.Column(visible=False):
            gr.Markdown(
                "**How does each training hyperparameter shape the loss curve?** "
                "The plots below are schematic (synthetic data, based on the "
                "standard behavior observed in the ML literature and on small "
                "transformers). Use them as a reference when choosing settings "
                "in the **Train** tab."
            )
            with gr.Row():
                with gr.Column(scale=1, min_width=220):
                    effect_radio = gr.Radio(
                        choices=effect_viz.PARAMS,
                        value=effect_viz.PARAMS[0],
                        label="Parameter",
                        interactive=True,
                        info="Select a hyperparameter to see its schematic effect on training and validation loss curves.",
                    )
                    effect_caption = gr.Markdown()
                with gr.Column(scale=2):
                    effect_img = gr.Image(
                        label="How it shapes training", type="filepath",
                        show_label=False, height=540,
                    )

            def on_effect_change(param):
                caption, img = render_effect_panel(param)
                return caption, img

            effect_radio.change(
                on_effect_change, inputs=effect_radio,
                outputs=[effect_caption, effect_img],
            )
            # Eager-render the first plot so the panel isn't blank on load
            demo.load(
                lambda: render_effect_panel(effect_viz.PARAMS[0]),
                inputs=None,
                outputs=[effect_caption, effect_img],
            )

        with gr.Column(visible=False) as _panel_reports:
            gr.Markdown(
                "### 16-slide visual walkthrough of a single forward pass\n\n"
                "Each slide reveals one stage of how the model processes your prompt — "
                "from raw text through tokenization, embeddings, Q/K/V projections, "
                "attention scores, causal masking, softmax, residual connections, "
                "the FFN layer, and finally the output logits and sampling.\n\n"
                "**How to use:** Enter a prompt and click *Render*. Adjust "
                "`layer` / `head` / `query_pos` to see how different parts of the "
                "model attend to different tokens. Slides 02–09 update when you "
                "change these controls.\n\n"
                "Aligned with Raschka's "
                "*Build a Large Language Model (From Scratch)* — see `REFERENCES.md`."
            )
            with gr.Row():
                with gr.Column(scale=1):
                    teach_prompt = gr.Textbox(
                        label="Prompt", value="The cat sat on the", lines=2,
                    )
                    teach_layer = gr.Slider(
                        0, n_layers - 1, 0, step=1, label="layer",
                    )
                    teach_head = gr.Slider(
                        0, n_heads - 1, 0, step=1, label="head",
                    )
                    teach_qpos = gr.Number(
                        value=-1, precision=0,
                        label="query_pos (-1 = last token)",
                    )
                    with gr.Row():
                        teach_temp = gr.Slider(
                            0.05, 2.0, 0.8, step=0.05, label="temp",
                        )
                        teach_topk = gr.Slider(
                            0, 200, 40, step=1, label="top_k",
                        )
                        teach_topp = gr.Slider(
                            0.05, 1.0, 0.9, step=0.05, label="top_p",
                        )
                    teach_btn = gr.Button("Render 16 slides", variant="primary")
                    with gr.Row():
                        teach_pin_btn = gr.Button("Pin for comparison", size="sm")
                        teach_dl_btn = gr.Button("Download PPTX", size="sm")
                    teach_dl_file = gr.File(label="Download", visible=False)
                    teach_status = gr.Markdown()
                with gr.Column(scale=2):
                    teach_pinned_label = gr.Markdown(visible=False)
                    teach_pinned_gallery = gr.Gallery(
                        label="Pinned (previous)", columns=2, height=350,
                        show_label=True, object_fit="contain", visible=False,
                    )
                    teach_gallery = gr.Gallery(
                        label="Current", columns=2, height=700,
                        show_label=True, object_fit="contain",
                    )
                    gr.Markdown(
                        "<small>"
                        "**Prompt** — text the model processes; short phrases (5–10 tokens) give the clearest slides.<br>"
                        f"**layer** — which transformer block (0–{n_layers-1}); early layers capture syntax, later layers capture semantics.<br>"
                        f"**head** — which attention head (0–{n_heads-1}); each head learns to attend to different token relationships.<br>"
                        "**query_pos** — which token is 'asking the question'; −1 = last token (most common for generation).<br>"
                        "**temp** — sampling temperature; lower = more deterministic, higher = more creative.<br>"
                        "**top_k** — only consider the top-k most likely tokens at each step; 0 = disabled.<br>"
                        "**top_p** — nucleus sampling: keep tokens until cumulative probability reaches top_p."
                        "</small>"
                    )
            def _render_and_track(*args):
                global _LAST_SLIDE_DIR
                import shutil
                old_dir = _LAST_SLIDE_DIR
                gallery, status, slide_dir = render_teach(*args)
                _LAST_SLIDE_DIR = slide_dir
                if old_dir and os.path.isdir(old_dir):
                    shutil.rmtree(old_dir, ignore_errors=True)
                return gallery, status

            teach_btn.click(
                _render_and_track,
                inputs=[teach_prompt, teach_layer, teach_head, teach_qpos,
                        teach_temp, teach_topk, teach_topp],
                outputs=[teach_gallery, teach_status],
            )

            def _pin_current(current_gallery, status_text):
                global _PINNED_GALLERY, _PINNED_LABEL
                if not current_gallery:
                    return gr.update(), gr.update(), gr.update()
                _PINNED_GALLERY = current_gallery
                _PINNED_LABEL = status_text or "Pinned"
                return (
                    gr.update(value=current_gallery, visible=True),
                    gr.update(value=f"**Pinned:** {_PINNED_LABEL}", visible=True),
                    gr.update(height=350),
                )

            teach_pin_btn.click(
                _pin_current,
                inputs=[teach_gallery, teach_status],
                outputs=[teach_pinned_gallery, teach_pinned_label, teach_gallery],
            )

            def _download_pptx(status_text):
                global _LAST_SLIDE_DIR
                if not _LAST_SLIDE_DIR:
                    return gr.update(visible=False)
                label = status_text or "Train Reports"
                path = export_slides_pptx(_LAST_SLIDE_DIR, label)
                if path and os.path.exists(path):
                    return gr.update(value=path, visible=True)
                return gr.update(visible=False)

            teach_dl_btn.click(
                _download_pptx,
                inputs=[teach_status],
                outputs=[teach_dl_file],
            )

            with gr.Accordion("Slide guide — what each of the 16 slides shows", open=True):
                gr.Markdown(
                    "**01 — Tokenization** · The prompt is encoded to UTF-8 bytes, then BPE merges reduce it to "
                    "fewer tokens. Three rows show: original characters → hex bytes → final token IDs. "
                    "Green tokens are BPE merges (learned); red tokens are raw single bytes. "
                    "The compression ratio tells you how efficiently the vocabulary represents your text.\n\n"
                    "**02 — Token embeddings** · Each token ID is looked up in the embedding table to produce "
                    "a `d_model`-dimensional vector (heatmap, first 32 dims shown). "
                    "The colour pattern reflects how similar tokens cluster together in embedding space — "
                    "words with similar meanings get similar vectors through training.\n\n"
                    "**03 — Q, K, V projections** · The embedding of each token is linearly projected three "
                    "times to produce Query, Key, and Value vectors for the selected head. "
                    "Q asks *what am I looking for?*, K announces *what do I contain?*, V carries *what I'll pass on*. "
                    "Each is `d_head`-dimensional; the heatmap uses the same colour scale for easy comparison.\n\n"
                    "**04 — Attention scores (raw)** · Pairwise dot products: `Q · Kᵀ / √d_head`. "
                    "Entry (i, j) measures how relevant token j is to token i. "
                    "The full matrix is shown *before* masking — notice the upper triangle is still visible. "
                    "High positive scores (dark blue) mean strong relevance.\n\n"
                    "**05 — Causal mask applied** · The upper triangle is set to −∞ so token i cannot "
                    "attend to any future token j > i. Grey cells are the masked positions. "
                    "This is what makes the model autoregressive: it can only use past context to "
                    "predict the next token.\n\n"
                    "**06 — Attention weights (softmax)** · After softmax, each row sums to 1 — it is now "
                    "a probability distribution over past tokens. The `*` marks the token each query "
                    "attends to most strongly. A uniform row means the head is uncertain; "
                    "a sharp row means it found a clear signal.\n\n"
                    "**07 — Weighted value sum** · For the selected query position (token i), the bar chart "
                    "shows how much weight it gave each past token. The heatmap below shows each V row "
                    "followed by the computed output = Σⱼ w_ij · V[j]. "
                    "This is the actual output of the attention head for that one token position.\n\n"
                    "**08 — All heads side by side** · Every attention head in the chosen layer at once. "
                    "Different heads typically specialise: one may copy the previous token, another may "
                    "track subject-verb agreement, another may link pronouns to their referents. "
                    "Diversity across heads is healthy — uniform heads suggest redundancy.\n\n"
                    "**09 — FFN delta (before / delta / after)** · After attention, the feed-forward network "
                    "processes each token independently. The three heatmaps show the hidden state entering "
                    "the FFN, the change the FFN adds (the *delta*), and the state after. "
                    "The delta panel shows which dimensions the FFN actively modifies — "
                    "this is where factual knowledge is believed to be stored.\n\n"
                    "**10 — Next-token distribution (top 20)** · The model's final prediction for what comes "
                    "next. Blue bars are the raw softmax probabilities; orange bars show what remains after "
                    "temperature scaling + top-k + top-p filtering. "
                    "A well-trained model concentrates probability on semantically plausible continuations.\n\n"
                    "**11 — Sampling rollout (10 steps)** · Ten decode steps shown as a grid. "
                    "Each column is one step; each row is one of the top-5 candidates with its probability. "
                    "The orange highlighted cell is the token that was actually chosen. "
                    "Watch how probabilities shift as context accumulates.\n\n"
                    "**12 — RoPE positional encoding** · Left: the cos/sin frequency tables — each column "
                    "oscillates at a different frequency, so every position gets a unique fingerprint. "
                    "Right: a unit vector rotated by increasing position shows how Q and K vectors are "
                    "literally rotated in 2D pairs. Low-frequency pairs rotate slowly (long-range), "
                    "high-frequency pairs rotate fast (short-range).\n\n"
                    "**13 — Why divide by √d_head?** · Without the scaling, Q·Kᵀ dot products grow in "
                    "magnitude with dimensionality, pushing softmax toward a one-hot distribution "
                    "(near-zero gradients = dead attention). "
                    "The grid compares unscaled vs scaled softmax at d = 8 / 64 / 256 / 1024 — "
                    "the top row saturates while the bottom row stays diffuse.\n\n"
                    "**14 — Temperature effect** · The same logits fed to softmax at T = 0.3 / 1.0 / 2.0. "
                    "T < 1 sharpens the distribution (model becomes confident, repetitive). "
                    "T = 1 is the raw distribution. "
                    "T > 1 flattens it (more creative but also more incoherent).\n\n"
                    "**15 — Greedy vs sampling** · Three decoded continuations from the same prompt: "
                    "greedy (always pick argmax, T ≈ 0), temperature-1 sampling (top-k=40, top-p=0.9), "
                    "and hot sampling (T=1.5, top-p=0.95). "
                    "Greedy output is often repetitive; hot sampling can be incoherent. "
                    "The middle setting is usually the best balance for story-like text.\n\n"
                    "**16 — Parameter breakdown** · Pie chart of where the 12M parameters live. "
                    "Token embedding / lm_head (weight-tied) and FFN layers dominate. "
                    "RMSNorm has almost none — it is cheap. "
                    "This explains why making the model wider (`d_model`) or deeper (more layers) "
                    "grows the FFN and QKV blocks fastest."
                )

        with gr.Column(visible=False) as _panel_attn:
            gr.Markdown(
                "### Attention heatmaps — what tokens attend to what\n\n"
                "Self-attention is the core mechanism that lets each token 'look at' "
                "every earlier token to decide what information to carry forward. "
                "This tab visualizes those attention weights directly.\n\n"
                "**Left (Single head):** One attention head's weight matrix after "
                "softmax + causal masking. Rows = query token (the one 'looking'), "
                "columns = key token (the one being 'looked at'). "
                "Bright = strong attention. The lower triangle is zero because "
                "causal masking prevents tokens from attending to future positions.\n\n"
                "**Right (Rollout):** Attention rollout across all layers "
                "(Abnar & Zuidema 2020) — approximates how information flows from "
                "input tokens to the final output across the *entire* model depth, "
                "not just one layer.\n\n"
                "**What the sliders do:**\n"
                "- **layer:** The model has {n_layers} stacked transformer blocks. "
                "Layer 0 (first) typically learns surface patterns — positional attention, "
                "punctuation, local word relationships. Higher layers learn increasingly "
                "abstract features — syntax, semantic roles, long-range dependencies. "
                "Changing the layer shows you a completely different stage of processing.\n"
                "- **head:** Each layer has {n_heads} parallel attention heads. "
                "Each head independently learns to focus on different types of relationships "
                "— one head might track subject-verb agreement, another might attend to "
                "the previous word, another to sentence boundaries. Changing the head shows "
                "you a different 'lens' on the same layer.\n\n"
                "**Try:** Compare layer 0 head 0 vs layer {n_layers_m1} head 0 to see "
                "how attention evolves from shallow to deep."
                .format(n_layers=n_layers, n_heads=n_heads, n_layers_m1=n_layers - 1)
            )
            with gr.Row():
                with gr.Column(scale=1):
                    attn_prompt = gr.Textbox(
                        label="Prompt", value="To be or not to be", lines=2,
                        info="Shorter prompts produce clearer heatmaps (5-10 tokens ideal).",
                    )
                    attn_layer = gr.Slider(
                        0, n_layers - 1, n_layers // 2, step=1, label="layer",
                        info=f"Layer 0 = surface patterns, layer {n_layers-1} = deep semantics. Changes the left heatmap.",
                    )
                    attn_head = gr.Slider(
                        0, n_heads - 1, 0, step=1, label="head",
                        info=f"Each of the {n_heads} heads learns different attention patterns. Changes the left heatmap.",
                    )
                    attn_btn = gr.Button("Render attention", variant="primary")
                    attn_status = gr.Markdown()
                with gr.Column(scale=2):
                    with gr.Row():
                        attn_head_img = gr.Image(label="Single head", type="filepath")
                        attn_rollout_img = gr.Image(label="Rollout (all layers)", type="filepath")
                    with gr.Row():
                        gr.Markdown(
                            "<small><b>Single head</b> — one attention head's weight matrix "
                            "after softmax + causal masking. "
                            "Rows are query tokens (the token 'asking'), columns are key tokens (the token 'answering'). "
                            "Brightness = attention weight (0–1, each row sums to 1). "
                            "The lower-left triangle only is filled because causal masking prevents any token from "
                            "attending to future positions. "
                            "A head that attends uniformly across past tokens is gathering broad context; "
                            "one that attends sharply to one token is copying or referencing it specifically. "
                            "Change <b>layer</b> and <b>head</b> to see how different parts of the model use attention differently.</small>"
                        )
                        gr.Markdown(
                            "<small><b>Rollout (all layers)</b> — attention rollout (Abnar & Zuidema 2020) propagates "
                            "attention weights through every layer simultaneously to approximate the true "
                            "information flow from input tokens to the final representation. "
                            "Unlike the single-head view, rollout accounts for residual connections and averages "
                            "across all heads at each layer. "
                            "A bright entry (i, j) means token j's input was a strong contributor to token i's "
                            "final hidden state. "
                            "This is more meaningful than any single head for understanding which input tokens "
                            "the model 'relied on' when building each position's representation.</small>"
                        )
            attn_btn.click(
                render_attention,
                inputs=[attn_prompt, attn_layer, attn_head],
                outputs=[attn_head_img, attn_rollout_img, attn_status],
            )

        with gr.Column(visible=False) as _panel_viz:
            gr.Markdown(
                "### Animated forward pass — watch data flow through the transformer\n\n"
                "Enter a prompt and click **Visualize** to see a synchronized animation of:\n"
                "- **Left:** Architecture diagram — each layer lights up as data flows through\n"
                "- **Center:** 6×6 attention grid — every head's attention pattern, revealed layer by layer\n"
                "- **Right:** Activation flow — hidden state magnitude at each layer\n\n"
                "Click any cell in the attention grid to see the full-size heatmap and tensor shapes. "
                "Use the playback controls (Play/Pause/Step) to control the animation speed."
            )
            viz_prompt = gr.Textbox(
                label="Prompt", value="To be or not to be",
                lines=1, max_lines=2,
                info="Text to process. Shorter prompts (5-10 tokens) produce clearer visualizations.",
            )
            viz_btn = gr.Button("Visualize forward pass", variant="primary")
            viz_html = gr.HTML(label="Visualization")

            viz_btn.click(
                render_visualization,
                inputs=[viz_prompt],
                outputs=[viz_html],
            )

        with gr.Column(visible=False) as _panel_gen:
            gr.Markdown(
                "### Interactive text generation\n\n"
                "The model predicts one token at a time, appending each prediction to "
                "the prompt and feeding it back in. The three sampling parameters "
                "(`temperature`, `top_k`, `top_p`) control *how* the next token is "
                "chosen from the model's probability distribution — they don't change "
                "the model itself, just how creative vs. deterministic the output is.\n\n"
                "**Tip:** Toggle the KV cache off to see the slower no-cache path "
                "(the reference implementation used in most tutorials). "
                "With cache on, each new token only computes attention against the "
                "cached K/V from previous tokens — O(T) instead of O(T²)."
            )
            with gr.Row():
                with gr.Column(scale=2):
                    gen_prompt = gr.Textbox(
                        label="Prompt", value="To be or not to be, ",
                        lines=2, max_lines=4,
                    )
                    with gr.Row():
                        gen_max = gr.Slider(
                            8, max_seq - 1, 100, step=1,
                            label="max_new_tokens",
                        )
                    with gr.Row():
                        gen_temp = gr.Slider(
                            0.05, 2.0, 0.8, step=0.05, label="temperature",
                        )
                        gen_topk = gr.Slider(
                            0, 200, 40, step=1, label="top_k",
                        )
                        gen_topp = gr.Slider(
                            0.05, 1.0, 0.9, step=0.05, label="top_p (nucleus)",
                        )
                    gen_cache = gr.Checkbox(
                        True, label="Use KV cache (generate_fast)",
                    )
                    gen_btn = gr.Button("Generate", variant="primary")
                with gr.Column(scale=3):
                    gen_out = gr.Textbox(
                        label="Output (streaming)", lines=16, show_copy_button=True,
                    )
                    gr.Markdown(
                        "<small>"
                        "**Prompt** — starting text the model continues from; it has only seen simple short stories during training.<br>"
                        "**max_new_tokens** — how many tokens to generate; longer = slower but more text; limited by the context window.<br>"
                        "**temperature** — scales logits before softmax; &lt;1 = sharper/more repetitive, &gt;1 = flatter/more random; 0.8 is a good default.<br>"
                        "**top_k** — only consider the top-k most probable tokens at each step; 0 = disabled (use all tokens).<br>"
                        "**top_p (nucleus)** — keep the smallest set of tokens whose cumulative probability exceeds top_p; adapts dynamically — broad when uncertain, narrow when confident.<br>"
                        "**KV cache** — ON: reuses past K/V vectors (fast, O(T) per step); OFF: recomputes full attention each step (slow, O(T²)) — toggle to see the educational comparison."
                        "</small>"
                    )
            gen_btn.click(
                generate_stream,
                inputs=[gen_prompt, gen_max, gen_temp, gen_topk, gen_topp, gen_cache],
                outputs=gen_out,
            )

        with gr.Column(visible=False) as _panel_bench:
            gr.Markdown(
                "### KV cache speedup — why caching matters\n\n"
                "Compares two generation strategies side by side:\n"
                "- **`generate()`** — no cache. Every new token recomputes attention "
                "over *all* previous tokens from scratch. Cost: O(T²) per step.\n"
                "- **`generate_fast()`** — KV cache. Stores each layer's K and V "
                "tensors from previous steps and only computes the new token's "
                "attention. Cost: O(T) per step.\n\n"
                "The speedup grows with sequence length — at 200 tokens the cached "
                "version can be 5-10x faster. This is how production LLMs (GPT, LLaMA, "
                "Claude) achieve fast inference."
            )
            with gr.Row():
                with gr.Column(scale=1):
                    bench_prompt = gr.Textbox(
                        label="Prompt", value="The ", lines=1,
                        info="Short prompt to seed generation. The benchmark measures the generation speed, not prompt processing.",
                    )
                    bench_n = gr.Slider(
                        20, max_seq - 10, 100, step=10,
                        label="tokens to generate",
                        info="More tokens = bigger speedup gap between cached and uncached. Try 200 for dramatic results.",
                    )
                    bench_btn = gr.Button("Run benchmark", variant="primary")
                    bench_summary = gr.Code(
                        label="Result", language=None, interactive=False,
                    )
                with gr.Column(scale=2):
                    bench_img = gr.Image(label="Throughput", type="filepath")
            bench_btn.click(
                run_bench,
                inputs=[bench_prompt, bench_n],
                outputs=[bench_img, bench_summary],
            )

        with gr.Column(visible=False) as _panel_kv:
            gr.Markdown(
                "### KV cache — equivalence proofs + length sweep\n\n"
                "Mirrors **§9** of the notebook. The Benchmark tab measures "
                "throughput on one prompt; this tab does the deeper analysis:\n\n"
                "1. **Single-step equivalence** — full forward on T tokens vs. "
                "(prefill on T-1) + (1-token decode with cache). Should match "
                "to ~`1e-7`.\n"
                "2. **Multi-step equivalence** — feed tokens one at a time "
                "through the cache, compare the final-position logits to a "
                "single full forward.\n"
                "3. **Length sweep** — generate `gen_len` new tokens at "
                "varying prompt lengths and compare `generate()` vs "
                "`generate_fast()`.\n\n"
                "If any equivalence test fails, attention or RoPE has a bug — "
                "this is the canary that `tests/test_model.py::test_kv_cache_*` "
                "guards in CI."
            )
            with gr.Row():
                with gr.Column(scale=1):
                    kv_T = gr.Slider(
                        4, max_seq, value=min(16, max_seq), step=1,
                        label="single-step: sequence length T",
                    )
                    kv_multi_T = gr.Slider(
                        4, max_seq, value=min(20, max_seq), step=1,
                        label="multi-step: tokens fed one-by-one",
                    )
                    kv_sweep = gr.Textbox(
                        label="length-sweep prompt lengths (comma-separated)",
                        value="5, 10, 20, 40",
                    )
                    kv_gen = gr.Slider(
                        10, max(20, max_seq - 50), value=min(50, max_seq // 2),
                        step=10, label="length-sweep: tokens to generate",
                    )
                    kv_btn = gr.Button("Run all 3 tests", variant="primary")
                    kv_single = gr.Code(label="Single-step", language=None,
                                        interactive=False)
                    kv_multi = gr.Code(label="Multi-step", language=None,
                                       interactive=False)
                with gr.Column(scale=2):
                    kv_sweep_img = gr.Image(label="Length-sweep timing",
                                            type="filepath")
                    kv_sweep_txt = gr.Code(label="Length-sweep numbers",
                                           language=None, interactive=False)
            kv_btn.click(
                run_kv_cache,
                inputs=[kv_T, kv_multi_T, kv_sweep, kv_gen],
                outputs=[kv_single, kv_multi, kv_sweep_img, kv_sweep_txt],
            )

        gr.Markdown(
            "\n---\n"
            "Built on PyTorch · GPT-style decoder · RMSNorm + RoPE + SwiGLU + "
            "Pre-Norm · weight-tied · combined QKV · bf16-capable. "
            "Code: `model.py` (800 lines)."
        )

        # ── Wire sidebar nav to panel visibility ──
        _MERMAID_DIAGRAM = """
flowchart TD
    subgraph SETUP["Setup  (bash run.sh setup)"]
        CORPUS["data/corpus.txt\\n~1.5M characters"]
        TRAIN_TOK["BPETokenizer.train()\\n256 base bytes → 4096 vocab"]
        TOK_JSON["tokenizer.json\\nvocab + merge rules"]
        CORPUS --> TRAIN_TOK --> TOK_JSON
    end

    subgraph TOKENISE["Tokenisation  (train.py)"]
        ENCODE["tokenizer.encode(corpus)\\nprogress every 5%"]
        CACHE["tokens_HASH.npy\\ndisk cache — skipped on 2nd run"]
        SPLIT["90 / 10 split\\ntrain_tokens · val_tokens"]
        DATASET["TextDataset\\nsliding windows, stride = seq_len/2"]
        LOADER["DataLoader\\nbatch → (input_ids, targets)\\ntargets = input shifted +1"]
        ENCODE --> CACHE --> SPLIT --> DATASET --> LOADER
    end

    subgraph MODEL["NanoLLM forward pass  (model.py)"]
        EMBTOK["token_emb\\ntoken ID → d_model vector"]
        subgraph BLOCK["× n_layers  TransformerBlock"]
            RMS1["RMSNorm"] --> ATTN["CausalSelfAttention\\nQKV · RoPE · scores/√d\\ncausal mask · softmax · out proj"]
            ATTN --> RES1["residual add"] --> RMS2["RMSNorm"]
            RMS2 --> FFN["SwiGLU FFN\\ngate · up → silu → down"]
            FFN --> RES2["residual add"]
        end
        LMHEAD["lm_head  (weight-tied to token_emb)\\nhidden → vocab logits"]
        LOSS["cross_entropy_loss\\nlogits vs targets"]
        EMBTOK --> BLOCK --> LMHEAD --> LOSS
    end

    subgraph TRAIN["Training loop  (train.py)"]
        AMP["AMP autocast  bf16/CUDA"]
        BACK["loss.backward()"]
        CLIP["clip_grad_norm"]
        OPT["AdamW step\\n+ cosine LR schedule"]
        CKPT["checkpoints/\\nepoch_NNN.pt · best.pt"]
        CURVE["loss_curve.png"]
        AMP --> BACK --> CLIP --> OPT --> CKPT
        OPT --> CURVE
    end

    subgraph INFER["Inference  (generate.py / model.py)"]
        GEN["generate()\\nno KV cache"]
        GENF["generate_fast()\\nKV cache — new token only"]
        SAMPLE["_sample_from_logits()\\ntemperature · top-k · top-p"]
        DECODE["tokenizer.decode()\\ntoken IDs → UTF-8 text"]
        GEN --> SAMPLE --> DECODE
        GENF --> SAMPLE
    end

    subgraph UI["Gradio UI  (app.py)"]
        T1["Tokenizer · Dataset · TransformerBlock"]
        T2["Train · Train Reports"]
        T3["Attention · Visualize"]
        T4["Generate · Benchmark · KV Cache"]
    end

    TOK_JSON --> TOKENISE
    LOADER --> MODEL
    LOSS --> TRAIN
    CKPT --> INFER
    CKPT --> UI
    TOK_JSON --> UI
    INFER --> UI
"""

        with gr.Column(visible=False) as _panel_arch:
            gr.Markdown("### End-to-end architecture flow")
            gr.HTML(f"""
<div style="overflow:auto; padding:16px;">
  <pre class="mermaid" style="background:transparent;">
{_MERMAID_DIAGRAM}
  </pre>
</div>
<script type="module">
  import mermaid from 'https://cdn.jsdelivr.net/npm/mermaid@11/dist/mermaid.esm.min.mjs';
  mermaid.initialize({{
    startOnLoad: true,
    theme: 'dark',
    flowchart: {{ curve: 'basis', padding: 20 }}
  }});
  mermaid.run();
</script>
""")

        _all_panels = [
            _panel_tokenizer, _panel_dataset, _panel_tb, _panel_train,
            _panel_finetune, _panel_reports, _panel_attn, _panel_viz, _panel_gen,
            _panel_bench, _panel_kv, _panel_arch,
        ]

        def _switch_panel(choice):
            return [gr.update(visible=(name == choice)) for name in _NAV_CHOICES]

        _nav.change(_switch_panel, inputs=_nav, outputs=_all_panels)

    return demo


# ═══════════════════════════════════════════════════════════════
# Entrypoint
# ═══════════════════════════════════════════════════════════════

def find_free_port(start_port: int, host: str = "127.0.0.1",
                   attempts: int = 20) -> int:
    """Return the first free port at or above `start_port`, probing by bind.

    Gradio's own retry only tries the exact port you pass. This helper walks
    the next N ports so `bash run.sh ui` Just Works even when a prior instance
    is still holding 7860.
    """
    for port in range(start_port, start_port + attempts):
        with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
            try:
                s.bind((host, port))
                return port
            except OSError:
                continue
    raise RuntimeError(
        f"No free port in {start_port}..{start_port + attempts - 1} on {host}"
    )


def main():
    parser = argparse.ArgumentParser(description="LLM101 - Gradio console")
    parser.add_argument("--checkpoint", default="checkpoints/best.pt")
    parser.add_argument("--port", type=int, default=7860,
                        help="Preferred starting port (will auto-fallback if busy)")
    parser.add_argument("--share", action="store_true",
                        help="Create a public URL via Gradio's share tunnel")
    parser.add_argument("--host", default="127.0.0.1",
                        help="Bind address (use 0.0.0.0 to expose on LAN)")
    args = parser.parse_args()

    print("LLM101 Console — starting...")
    _load_model(args.checkpoint)
    print(f"  {_STATUS}")

    # Auto-fallback to the next free port so a running instance on 7860 doesn't
    # crash a fresh `bash run.sh ui` call.
    port = find_free_port(args.port, host=args.host)
    if port != args.port:
        print(f"  Port {args.port} is busy — using {port} instead")

    demo = build_ui()
    # show_api=False works around a gradio-client schema-introspection bug
    # (TypeError on schema traversal in gradio 4.44.x). The UI still works,
    # we just skip the auto-generated API docs page.
    demo.queue().launch(
        server_name=args.host,
        server_port=port,
        share=args.share,
        show_error=True,
        inbrowser=True,
        show_api=False,
    )


if __name__ == "__main__":
    main()
